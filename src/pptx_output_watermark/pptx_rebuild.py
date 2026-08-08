from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from .models import WatermarkOptions
from .watermarking import slide_emu_to_overlay_pixels, write_watermark_overlay_image


def build_image_pptx_from_images(
    image_paths: list[Path],
    output_path: Path,
    *,
    slide_width_emu: int | None = None,
    slide_height_emu: int | None = None,
) -> Path:
    if not image_paths:
        raise RuntimeError("No images available to build image-based PPTX.")

    prs = Presentation()
    if slide_width_emu is not None and slide_height_emu is not None:
        prs.slide_width = Emu(int(slide_width_emu))
        prs.slide_height = Emu(int(slide_height_emu))
    else:
        with Image.open(image_paths[0]) as first_image:
            width_px, height_px = first_image.size
        prs.slide_width = Emu(int(width_px * 9525))
        prs.slide_height = Emu(int(height_px * 9525))
    blank_layout = prs.slide_layouts[6]
    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(output_path)
    return output_path


def add_watermark_overlay_to_pptx(
    input_pptx: Path,
    output_pptx: Path,
    options: WatermarkOptions,
) -> Path:
    prs = Presentation(str(input_pptx))
    overlay_png = write_watermark_overlay_image(
        slide_emu_to_overlay_pixels(prs.slide_width, prs.slide_height),
        options,
    )
    try:
        for slide in prs.slides:
            slide.shapes.add_picture(
                str(overlay_png),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        prs.save(str(output_pptx))
    finally:
        try:
            overlay_png.unlink()
        except FileNotFoundError:
            pass
    return output_pptx


def copy_pptx(input_pptx: Path, output_pptx: Path) -> Path:
    shutil.copy2(input_pptx, output_pptx)
    return output_pptx
