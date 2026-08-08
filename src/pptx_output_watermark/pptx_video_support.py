from __future__ import annotations

import json
import mimetypes
import os
import posixpath
import re
import shutil
import subprocess
import sys
import time
from collections.abc import Callable
from dataclasses import dataclass, field
from functools import lru_cache
from pathlib import Path
from tempfile import mkdtemp, mkstemp
from zipfile import ZipFile
from xml.etree import ElementTree as ET

from PIL import Image
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx import Presentation
from pptx.util import Emu

from pptx_video_compactor import (
    PROFILE_QUALITY_RULES,
    VideoAsset as CompactVideoAsset,
    apply_video_bitrate_bounds,
    height_bucket_for_source,
    load_runtime_config,
    profile_target_fps,
    recommend_audio_kbps,
    rendition_for_asset,
    scale_dims,
)

from .ffmpeg_runtime import ensure_binary, run_binary
from .models import WatermarkOptions
from .process_utils import (
    finish_process,
    hidden_console_kwargs,
    kill_process,
    start_process,
    subprocess_text_kwargs,
)
from .watermarking import apply_watermark_to_image, write_watermark_overlay_image

VIDEO_EXTENSIONS = {
    ".mp4",
    ".m4v",
    ".mov",
    ".wmv",
    ".avi",
    ".asf",
    ".mpg",
    ".mpeg",
    ".webm",
    ".mkv",
    ".ts",
}
GPU_ENCODERS = ("h264_videotoolbox", "h264_nvenc", "h264_qsv", "h264_amf", "h264_mf")

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}
REL_NS = {"pr": "http://schemas.openxmlformats.org/package/2006/relationships"}
R_NS = f"{{{NS['r']}}}"
PRESENTATION_PATH = "ppt/presentation.xml"
PRESENTATION_RELS_PATH = "ppt/_rels/presentation.xml.rels"


@dataclass(slots=True)
class VideoOccurrence:
    slide_number: int
    occurrence_index: int
    slide_path: str
    shape_id: int
    relationship_ids: list[str]
    shape_name: str
    media_path: str
    x: int
    y: int
    cx: int
    cy: int
    rotation: str | None = None
    flip_h: str | None = None
    flip_v: str | None = None
    crop: dict[str, str] = field(default_factory=dict)


@dataclass(slots=True)
class VideoAsset:
    media_path: str
    occurrences: list[VideoOccurrence] = field(default_factory=list)
    extracted_path: Path | None = None
    prepared_path: Path | None = None
    poster_path: Path | None = None
    mime_type: str = "video/mp4"


@dataclass(slots=True)
class VideoSidecarInfo:
    directory: Path
    manifest_path: Path


@dataclass(slots=True)
class VideoTranscodeProfile:
    frame_rate: str | None = None
    video_bitrate: int | None = None
    audio_bitrate: int | None = None
    duration_seconds: float | None = None
    width: int = 0
    height: int = 0
    video_filter: str | None = None


Logger = Callable[[str], None]
ProgressLogger = Callable[[int], None]


def _slide_sort_key(path: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", path)
    return int(match.group(1)) if match else 0


def _slide_fallback_paths(zf: ZipFile) -> list[str]:
    return sorted(
        (
            name
            for name in zf.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ),
        key=_slide_sort_key,
    )


def _presentation_ordered_slide_paths(zf: ZipFile) -> list[str]:
    if (
        PRESENTATION_PATH not in zf.namelist()
        or PRESENTATION_RELS_PATH not in zf.namelist()
    ):
        return _slide_fallback_paths(zf)

    try:
        presentation = ET.fromstring(zf.read(PRESENTATION_PATH))
        rels_root = ET.fromstring(zf.read(PRESENTATION_RELS_PATH))
    except ET.ParseError:
        return _slide_fallback_paths(zf)

    relationships: dict[str, str] = {}
    for rel in rels_root.findall("pr:Relationship", REL_NS):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if not rel_id or not target:
            continue
        relationships[rel_id] = target

    slide_paths: list[str] = []
    for slide_id in presentation.findall(".//p:sldIdLst/p:sldId", NS):
        rel_id = slide_id.attrib.get(f"{R_NS}id")
        target = relationships.get(rel_id or "")
        if not target:
            continue
        slide_path = (
            target
            if target.startswith("ppt/")
            else posixpath.normpath(posixpath.join("ppt", target))
        )
        if slide_path in zf.namelist():
            slide_paths.append(slide_path)

    return slide_paths or _slide_fallback_paths(zf)


def _is_hidden_slide(zf: ZipFile, slide_path: str) -> bool:
    try:
        root = ET.fromstring(zf.read(slide_path))
    except ET.ParseError:
        return False
    return root.attrib.get("show") == "0"


def visible_slide_number_map(pptx_path: Path) -> dict[int, int]:
    """Return source slide number -> exported visible slide number."""
    mapping: dict[int, int] = {}
    visible_number = 0
    with ZipFile(pptx_path) as zf:
        for source_number, slide_path in enumerate(
            _presentation_ordered_slide_paths(zf), start=1
        ):
            if _is_hidden_slide(zf, slide_path):
                continue
            visible_number += 1
            mapping[source_number] = visible_number
    return mapping


def slide_visibility_counts(pptx_path: Path) -> tuple[int, int]:
    with ZipFile(pptx_path) as zf:
        slide_paths = _presentation_ordered_slide_paths(zf)
        hidden_count = sum(
            1 for slide_path in slide_paths if _is_hidden_slide(zf, slide_path)
        )
        return len(slide_paths), hidden_count


def _relationships_path_for(slide_path: str) -> str:
    base_dir = posixpath.dirname(slide_path)
    filename = posixpath.basename(slide_path)
    return posixpath.join(base_dir, "_rels", f"{filename}.rels")


def _resolve_zip_target(source_path: str, target: str) -> str:
    source_dir = posixpath.dirname(source_path)
    return posixpath.normpath(posixpath.join(source_dir, target)).lstrip("/")


def _parse_relationships(zf: ZipFile, slide_path: str) -> dict[str, str]:
    rels_path = _relationships_path_for(slide_path)
    if rels_path not in zf.namelist():
        return {}
    root = ET.fromstring(zf.read(rels_path))
    relationships: dict[str, str] = {}
    for rel in root.findall("pr:Relationship", REL_NS):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if not rel_id or not target or rel.attrib.get("TargetMode") == "External":
            continue
        relationships[rel_id] = _resolve_zip_target(slide_path, target)
    return relationships


def _video_rel_ids(pic: ET.Element) -> list[str]:
    nv_pic = pic.find("./p:nvPicPr", NS)
    if nv_pic is None:
        return []
    nv_pr = nv_pic.find("./p:nvPr", NS)
    if nv_pr is None:
        return []
    rel_ids: list[str] = []
    video_file = nv_pr.find("a:videoFile", NS)
    if video_file is not None:
        rel_id = video_file.attrib.get(f"{R_NS}link")
        if rel_id:
            rel_ids.append(rel_id)
    media_ref = nv_pr.find(".//p14:media", NS)
    if media_ref is not None:
        rel_id = media_ref.attrib.get(f"{R_NS}embed")
        if rel_id:
            rel_ids.append(rel_id)
    return rel_ids


def _shape_crop_attrs(pic: ET.Element) -> dict[str, str]:
    src_rect = pic.find("./p:blipFill/a:srcRect", NS)
    if src_rect is None:
        return {}
    return {
        key: value
        for key, value in src_rect.attrib.items()
        if key in {"l", "t", "r", "b"} and value
    }


def scan_embedded_videos(
    pptx_path: Path,
    *,
    slide_number_map: dict[int, int] | None = None,
) -> dict[str, VideoAsset]:
    assets: dict[str, VideoAsset] = {}
    with ZipFile(pptx_path) as zf:
        slide_paths = _presentation_ordered_slide_paths(zf)
        for source_slide_number, slide_path in enumerate(slide_paths, start=1):
            if slide_number_map is not None:
                slide_number = slide_number_map.get(source_slide_number)
                if slide_number is None:
                    continue
            else:
                slide_number = source_slide_number
            relationships = _parse_relationships(zf, slide_path)
            root = ET.fromstring(zf.read(slide_path))
            occurrence_index = 0
            for pic in root.findall(".//p:pic", NS):
                rel_ids = _video_rel_ids(pic)
                if not rel_ids:
                    continue
                media_path = ""
                for rel_id in rel_ids:
                    candidate = relationships.get(rel_id, "")
                    if Path(candidate).suffix.lower() in VIDEO_EXTENSIONS:
                        media_path = candidate
                        break
                if not media_path:
                    continue

                occurrence_index += 1
                nv_pic = pic.find("./p:nvPicPr", NS)
                c_nv_pr = None if nv_pic is None else nv_pic.find("./p:cNvPr", NS)
                xfrm = pic.find("./p:spPr/a:xfrm", NS)
                if xfrm is None:
                    continue
                ext = xfrm.find("a:ext", NS)
                off = xfrm.find("a:off", NS)
                if ext is None or off is None:
                    continue

                occurrence = VideoOccurrence(
                    slide_number=slide_number,
                    occurrence_index=occurrence_index,
                    slide_path=slide_path,
                    shape_id=int(c_nv_pr.attrib.get("id", "0"))
                    if c_nv_pr is not None
                    else 0,
                    relationship_ids=rel_ids,
                    shape_name=(
                        c_nv_pr.attrib.get(
                            "name", f"slide{slide_number}_video_{occurrence_index}"
                        )
                        if c_nv_pr is not None
                        else f"slide{slide_number}_video_{occurrence_index}"
                    ),
                    media_path=media_path,
                    x=int(off.attrib["x"]),
                    y=int(off.attrib["y"]),
                    cx=int(ext.attrib["cx"]),
                    cy=int(ext.attrib["cy"]),
                    rotation=xfrm.attrib.get("rot"),
                    flip_h=xfrm.attrib.get("flipH"),
                    flip_v=xfrm.attrib.get("flipV"),
                    crop=_shape_crop_attrs(pic),
                )
                asset = assets.setdefault(media_path, VideoAsset(media_path=media_path))
                asset.occurrences.append(occurrence)
                guessed_mime = mimetypes.guess_type(media_path)[0]
                if guessed_mime:
                    asset.mime_type = guessed_mime
    return assets


def extract_embedded_videos(
    pptx_path: Path, assets: dict[str, VideoAsset], output_dir: Path
) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    with ZipFile(pptx_path) as zf:
        for index, asset in enumerate(
            sorted(assets.values(), key=_asset_sort_key), start=1
        ):
            suffix = Path(asset.media_path).suffix.lower() or ".bin"
            target = output_dir / f"video_{index:03d}{suffix}"
            with zf.open(asset.media_path, "r") as source, target.open("wb") as output:
                shutil.copyfileobj(source, output, length=1024 * 1024)
            asset.extracted_path = target


def _asset_sort_key(asset: VideoAsset) -> tuple[int, int, str]:
    if not asset.occurrences:
        return (10**9, 10**9, asset.media_path)
    first = min(
        asset.occurrences, key=lambda item: (item.slide_number, item.occurrence_index)
    )
    return (first.slide_number, first.occurrence_index, asset.media_path)


def _duplicate_media_descriptions(assets: dict[str, VideoAsset]) -> list[str]:
    descriptions: list[str] = []
    for asset in sorted(assets.values(), key=_asset_sort_key):
        if len(asset.occurrences) <= 1:
            continue
        pages = ", ".join(str(item.slide_number) for item in asset.occurrences)
        descriptions.append(f"{Path(asset.media_path).name}: pages {pages}")
    return descriptions


def _extract_poster_frame(source_path: Path, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    run_binary(
        [
            "ffmpeg",
            "-y",
            "-v",
            "error",
            "-nostats",
            "-i",
            str(source_path),
            "-frames:v",
            "1",
            "-update",
            "1",
            str(output_path),
        ]
    )
    return output_path


def extract_video_poster_frame(source_path: Path, output_path: Path) -> Path:
    return _extract_poster_frame(source_path, output_path)


def _positive_int(value: object) -> int | None:
    try:
        parsed = int(str(value))
    except (TypeError, ValueError):
        return None
    return parsed if parsed > 0 else None


def _positive_float(value: object) -> float | None:
    try:
        parsed = float(str(value))
    except (TypeError, ValueError):
        return None
    return parsed if parsed > 0 else None


def _ffmpeg_progress_seconds(key: str, value: str) -> float | None:
    """Normalize FFmpeg's current and legacy progress timestamp fields."""
    if key == "out_time":
        parts = value.split(":")
        if len(parts) != 3:
            return None
        try:
            return max(0.0, int(parts[0]) * 3600 + int(parts[1]) * 60 + float(parts[2]))
        except ValueError:
            return None
    timestamp = _positive_float(value)
    if timestamp is None:
        return None
    return timestamp / 1_000_000.0 if key in {"out_time_us", "out_time_ms"} else None


def _valid_frame_rate(value: object) -> str | None:
    if not isinstance(value, str) or "/" not in value:
        return None
    numerator_text, denominator_text = value.split("/", 1)
    numerator = _positive_int(numerator_text)
    denominator = _positive_int(denominator_text)
    if numerator is None or denominator is None:
        return None
    # Ignore pathological probe values. Normal presentation videos are far below this.
    if numerator / denominator > 240:
        return None
    return f"{numerator}/{denominator}"


def _probe_transcode_profile(source_path: Path) -> VideoTranscodeProfile:
    try:
        completed = run_binary(
            [
                "ffprobe",
                "-v",
                "error",
                "-print_format",
                "json",
                "-show_streams",
                "-show_format",
                str(source_path),
            ],
            capture=True,
        )
    except (FileNotFoundError, subprocess.CalledProcessError, json.JSONDecodeError):
        return VideoTranscodeProfile()

    try:
        payload = json.loads(completed.stdout or "{}")
    except json.JSONDecodeError:
        return VideoTranscodeProfile()

    streams = payload.get("streams") or []
    format_info = payload.get("format") or {}
    video_stream = next(
        (stream for stream in streams if stream.get("codec_type") == "video"), {}
    )
    audio_stream = next(
        (stream for stream in streams if stream.get("codec_type") == "audio"), {}
    )

    frame_rate = _valid_frame_rate(video_stream.get("avg_frame_rate"))
    if frame_rate is None:
        frame_rate = _valid_frame_rate(video_stream.get("r_frame_rate"))

    profile = VideoTranscodeProfile(
        frame_rate=frame_rate,
        video_bitrate=(
            _positive_int(video_stream.get("bit_rate"))
            or _positive_int(format_info.get("bit_rate"))
        ),
        audio_bitrate=_positive_int(audio_stream.get("bit_rate")),
        duration_seconds=(
            _positive_float(video_stream.get("duration"))
            or _positive_float(format_info.get("duration"))
        ),
        width=_positive_int(video_stream.get("width")) or 0,
        height=_positive_int(video_stream.get("height")) or 0,
    )
    return profile


def _apply_compactor_video_profile(
    source_path: Path,
    profile: VideoTranscodeProfile,
    quality_profile: str,
) -> VideoTranscodeProfile:
    if profile.width <= 0 or profile.height <= 0:
        return profile

    try:
        internal_profile = (
            quality_profile if quality_profile in PROFILE_QUALITY_RULES else "high"
        )
        config = load_runtime_config(None, internal_profile)
        rules = PROFILE_QUALITY_RULES[internal_profile]
        asset = CompactVideoAsset(
            media_path=source_path.name, zip_size=source_path.stat().st_size
        )
        asset.width = profile.width
        asset.height = profile.height
        asset.duration_sec = max(0.1, profile.duration_seconds or 0.1)
        asset.has_audio = profile.audio_bitrate is not None
        asset.original_audio_kbps = max(0, (profile.audio_bitrate or 0) // 1000)
        probed_video_kbps = max(0, (profile.video_bitrate or 0) // 1000)
        estimated_total_kbps = max(
            1, int(asset.zip_size * 8 / asset.duration_sec / 1000)
        )
        asset.original_total_kbps = max(
            probed_video_kbps + asset.original_audio_kbps,
            estimated_total_kbps,
            1,
        )
        asset.original_video_kbps = max(
            1,
            probed_video_kbps
            or max(asset.original_total_kbps - asset.original_audio_kbps, 1),
        )
        asset.selected_height = height_bucket_for_source(
            asset,
            min(
                int(rules["max_height"]), int(config.render_limits["max_output_height"])
            ),
        )
        rendition = rendition_for_asset(asset, config)
        asset.audio_kbps = recommend_audio_kbps(asset, 1.0, config)
        apply_video_bitrate_bounds(asset, rendition)
        dynamic_range = max(0, asset.max_video_kbps - asset.min_video_kbps)
        asset.target_video_kbps = asset.min_video_kbps + int(
            dynamic_range * float(rules["bitrate_bias"])
        )
        asset.target_video_kbps = min(asset.target_video_kbps, asset.max_video_kbps)
        asset.target_video_kbps = max(asset.target_video_kbps, asset.min_video_kbps)
        asset.target_fps = profile_target_fps(asset, internal_profile)
        scaled_w, scaled_h = scale_dims(asset, config)
    except Exception:
        return profile

    profile.video_bitrate = max(1, asset.target_video_kbps) * 1000
    if asset.audio_kbps > 0:
        profile.audio_bitrate = asset.audio_kbps * 1000
    if (scaled_w, scaled_h) != (profile.width, profile.height):
        profile.video_filter = f"scale={scaled_w}:{scaled_h}:flags=lanczos"
    return profile


@lru_cache(maxsize=1)
def available_encoder_names() -> set[str]:
    try:
        result = run_binary(["ffmpeg", "-hide_banner", "-encoders"], capture=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        return set()
    names: set[str] = set()
    output = (result.stdout or "") + (result.stderr or "")
    for line in output.splitlines():
        parts = line.split()
        if len(parts) >= 2 and parts[0].startswith("V"):
            names.add(parts[1])
    return names


def gpu_encoder_priority() -> tuple[str, ...]:
    if sys.platform == "darwin":
        return ("h264_videotoolbox",)
    if os.name == "nt":
        return ("h264_nvenc", "h264_qsv", "h264_amf", "h264_mf")
    return GPU_ENCODERS


def probe_gpu_encoder(encoder: str) -> bool:
    cmd = [
        "ffmpeg",
        "-hide_banner",
        "-loglevel",
        "error",
        "-f",
        "lavfi",
        "-i",
        "testsrc2=size=128x128:rate=24:duration=0.2",
        "-frames:v",
        "1",
        "-c:v",
        encoder,
        "-pix_fmt",
        "yuv420p",
        "-f",
        "null",
        "NUL" if os.name == "nt" else "/dev/null",
    ]
    if encoder == "h264_videotoolbox":
        cmd[-5:-5] = ["-profile:v", "main"]
    try:
        run_binary(cmd, capture=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        return False
    return True


@lru_cache(maxsize=1)
def usable_gpu_encoder_names() -> tuple[str, ...]:
    available = available_encoder_names()
    usable: list[str] = []
    for encoder in gpu_encoder_priority():
        if encoder in available and probe_gpu_encoder(encoder):
            usable.append(encoder)
    return tuple(usable)


def _video_quality_args(profile: VideoTranscodeProfile) -> list[str]:
    args: list[str] = []
    if profile.video_bitrate:
        args.extend(["-b:v", str(profile.video_bitrate)])
    else:
        args.extend(["-crf", "18"])
    return args


def _gpu_video_quality_args(profile: VideoTranscodeProfile) -> list[str]:
    if not profile.video_bitrate:
        return []
    return [
        "-b:v",
        str(profile.video_bitrate),
        "-maxrate",
        str(max(int(profile.video_bitrate * 1.25), profile.video_bitrate + 100_000)),
        "-bufsize",
        str(max(int(profile.video_bitrate * 2), profile.video_bitrate + 200_000)),
    ]


def _build_watermarked_video_command(
    source_path: Path,
    overlay_path: Path,
    output_path: Path,
    profile: VideoTranscodeProfile,
    *,
    encoder: str = "libx264",
    copy_audio: bool,
) -> list[str]:
    audio_args = ["-c:a", "copy"]
    if not copy_audio:
        audio_args = ["-c:a", "aac", "-b:a", str(profile.audio_bitrate or 160_000)]
    video_args = [
        "-c:v",
        encoder,
        "-pix_fmt",
        "yuv420p",
    ]
    if profile.frame_rate:
        video_args.extend(["-r", profile.frame_rate])
    if encoder == "libx264":
        video_args.extend(["-preset", "medium", *_video_quality_args(profile)])
    else:
        if encoder == "h264_videotoolbox":
            video_args.extend(["-profile:v", "main"])
        video_args.extend(_gpu_video_quality_args(profile))
    video_filter = "[0:v:0][1:v:0]overlay=0:0:format=auto"
    if profile.video_filter:
        video_filter = f"{video_filter},{profile.video_filter}"

    command = [
        "ffmpeg",
        "-y",
        "-i",
        str(source_path),
        "-i",
        str(overlay_path),
        "-filter_complex",
        f"{video_filter}[wmv]",
        "-map",
        "[wmv]",
        "-map",
        "0:a?",
        "-map_metadata",
        "0",
        "-map_chapters",
        "0",
        *video_args,
        *audio_args,
        "-movflags",
        "+faststart",
        str(output_path),
    ]
    return command


def _run_ffmpeg_encode_with_progress(
    command: list[str],
    *,
    duration_seconds: float | None,
    progress_logger: ProgressLogger | None,
) -> None:
    if not progress_logger or not duration_seconds:
        run_binary(command)
        return

    resolved_cmd = list(command)
    resolved_cmd[0] = ensure_binary(resolved_cmd[0])
    resolved_cmd[1:1] = ["-v", "error", "-nostats", "-progress", "pipe:1"]
    process = start_process(
        resolved_cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        **subprocess_text_kwargs(),
        **hidden_console_kwargs(),
    )
    last_bucket = 0
    try:
        if process.stdout is not None:
            for raw_line in process.stdout:
                key, _, value = raw_line.strip().partition("=")
                seconds = _ffmpeg_progress_seconds(key, value)
                if seconds is not None:
                    percent = int(
                        min(
                            99,
                            max(0, seconds / duration_seconds * 100),
                        )
                    )
                    bucket = (percent // 25) * 25
                    if bucket in {25, 50, 75} and bucket > last_bucket:
                        last_bucket = bucket
                        progress_logger(bucket)
                elif key == "progress" and value == "end":
                    progress_logger(100)
        stderr = process.stderr.read() if process.stderr is not None else ""
        returncode = process.wait()
    except Exception:
        kill_process(process)
        process.wait()
        raise
    finally:
        finish_process(process)
    if returncode:
        raise subprocess.CalledProcessError(returncode, resolved_cmd, stderr=stderr)


def _encode_watermarked_video(
    source_path: Path,
    overlay_path: Path,
    output_path: Path,
    *,
    encoder_mode: str = "auto",
    quality_profile: str = "high",
    logger: Logger | None = None,
    progress_logger: ProgressLogger | None = None,
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    profile = _apply_compactor_video_profile(
        source_path,
        _probe_transcode_profile(source_path),
        quality_profile,
    )
    if encoder_mode in {"auto", "gpu"}:
        gpu_encoders = usable_gpu_encoder_names()
        if not gpu_encoders:
            _log(logger, "No usable GPU H.264 encoder found; falling back to CPU x264.")
        elif not profile.video_bitrate:
            _log(
                logger, "Source video bitrate is unavailable; falling back to CPU x264."
            )
        else:
            for gpu_encoder in gpu_encoders:
                try:
                    _log(logger, f"Video encoder: GPU ({gpu_encoder})")
                    _run_ffmpeg_encode_with_progress(
                        _build_watermarked_video_command(
                            source_path,
                            overlay_path,
                            output_path,
                            profile,
                            encoder=gpu_encoder,
                            copy_audio=True,
                        ),
                        duration_seconds=profile.duration_seconds,
                        progress_logger=progress_logger,
                    )
                    return output_path
                except subprocess.CalledProcessError as exc:
                    try:
                        output_path.unlink()
                    except FileNotFoundError:
                        pass
                    _log(
                        logger,
                        f"GPU encoder {gpu_encoder} failed; trying next encoder.",
                    )
                    if encoder_mode == "gpu":
                        _log(logger, f"GPU failure detail: returncode={exc.returncode}")
            _log(
                logger,
                "All usable GPU encoders failed for this video; falling back to CPU x264.",
            )

    _log(logger, "Video encoder: CPU (libx264)")
    try:
        _run_ffmpeg_encode_with_progress(
            _build_watermarked_video_command(
                source_path,
                overlay_path,
                output_path,
                profile,
                encoder="libx264",
                copy_audio=True,
            ),
            duration_seconds=profile.duration_seconds,
            progress_logger=progress_logger,
        )
    except subprocess.CalledProcessError:
        try:
            output_path.unlink()
        except FileNotFoundError:
            pass
        _run_ffmpeg_encode_with_progress(
            _build_watermarked_video_command(
                source_path,
                overlay_path,
                output_path,
                profile,
                encoder="libx264",
                copy_audio=False,
            ),
            duration_seconds=profile.duration_seconds,
            progress_logger=progress_logger,
        )
    return output_path


def watermark_video_file(
    source_path: Path,
    output_path: Path,
    watermark: WatermarkOptions,
    *,
    encoder_mode: str = "auto",
    quality_profile: str = "high",
    logger: Logger | None = None,
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    profile = _probe_transcode_profile(source_path)
    if profile.width <= 0 or profile.height <= 0:
        raise RuntimeError(
            f"Unable to determine video dimensions for watermarking: {source_path}"
        )
    overlay_path = write_watermark_overlay_image(
        (profile.width, profile.height),
        watermark,
    )
    try:
        return _encode_watermarked_video(
            source_path,
            overlay_path,
            output_path,
            encoder_mode=encoder_mode,
            quality_profile=quality_profile,
            logger=logger,
        )
    finally:
        try:
            overlay_path.unlink()
        except FileNotFoundError:
            pass


def _log(logger: Logger | None, message: str) -> None:
    if logger is not None:
        logger(message)


def replace_video_media_in_pptx(
    pptx_path: Path,
    replacements: dict[str, Path],
    *,
    output_path: Path | None = None,
    logger: Logger | None = None,
) -> Path:
    """Replace video media bytes inside a PPTX zip in-place (zip copy).

    Unlike the image-PPTX rebuild path, this preserves the original slide XML
    (including p:timing animation trees, playback settings, triggers, audio
    shapes, and z-order) because it only swaps media file bytes — exactly like
    the compression pipeline. Used by the editable-PPTX + watermark path to
    watermarked videos without losing any playback fidelity.

    ``replacements`` maps zip-internal media paths (e.g. ``ppt/media/media1.mp4``)
    to local file paths of the watermarked replacements. ``output_path``
    defaults to overwriting ``pptx_path`` via a temp file + atomic move.
    """
    target = output_path or pptx_path
    target.parent.mkdir(parents=True, exist_ok=True)
    tmp_path = target.with_name(
        f".{target.stem}.tmp-{os.getpid()}-{int(time.time() * 1000)}{target.suffix}"
    )
    try:
        with (
            ZipFile(pptx_path, "r") as zin,
            ZipFile(tmp_path, "w", allowZip64=True) as zout,
        ):
            for info in zin.infolist():
                if info.filename in replacements:
                    replacement = replacements[info.filename]
                    with replacement.open("rb") as source:
                        zout.writestr(info, source.read())
                    _log(
                        logger,
                        f"Replaced media: {info.filename} <- {replacement.name}",
                    )
                else:
                    with zin.open(info, "r") as source:
                        zout.writestr(info, source.read())
        shutil.move(str(tmp_path), str(target))
        return target
    except Exception:
        try:
            tmp_path.unlink()
        except FileNotFoundError:
            pass
        raise


def watermark_videos_in_editable_pptx(
    pptx_path: Path,
    output_path: Path,
    watermark: WatermarkOptions,
    *,
    encoder_mode: str = "auto",
    quality_profile: str = "high",
    logger: Logger | None = None,
) -> Path:
    """Watermark all embedded videos in an editable PPTX via media replacement.

    Scans the PPTX for embedded videos, extracts each, watermarks it
    (re-encoding via ffmpeg), and replaces the original media bytes in the
    zip — preserving all slide XML, timing, playback settings, and audio.
    """
    assets = scan_embedded_videos(pptx_path)
    if not assets:
        _log(logger, "No embedded videos found to watermark")
        return _copy_pptx(pptx_path, output_path)

    work_dir = Path(mkdtemp(prefix="pptx_watermark_video_"))
    try:
        extract_embedded_videos(pptx_path, assets, work_dir)
        replacements: dict[str, Path] = {}
        sorted_assets = sorted(assets.values(), key=_asset_sort_key)
        total = len(sorted_assets)
        for index, asset in enumerate(sorted_assets, start=1):
            if asset.extracted_path is None:
                continue
            video_name = Path(asset.media_path).name
            _log(logger, f"Video {index}/{total}: watermarking ({video_name})")
            output_suffix = ".mp4"
            prepared = work_dir / f"video_{index:03d}_wm{output_suffix}"
            watermark_video_file(
                asset.extracted_path,
                prepared,
                watermark,
                encoder_mode=encoder_mode,
                quality_profile=quality_profile,
                logger=logger,
            )
            replacements[asset.media_path] = prepared
            asset.prepared_path = prepared
        _log(logger, f"Replacing {len(replacements)} video media file(s) in PPTX")
        return replace_video_media_in_pptx(
            pptx_path, replacements, output_path=output_path, logger=logger
        )
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


def _copy_pptx(source: Path, target: Path) -> Path:
    target.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, target)
    return target


def prepare_videos_for_image_pptx(
    pptx_path: Path,
    work_dir: Path,
    watermark: WatermarkOptions,
    *,
    encoder_mode: str = "auto",
    slide_number_map: dict[int, int] | None = None,
    logger: Logger | None = None,
) -> dict[str, VideoAsset]:
    raw_assets = scan_embedded_videos(pptx_path)
    assets = scan_embedded_videos(pptx_path, slide_number_map=slide_number_map)
    if not assets:
        return assets

    raw_placement_count = sum(len(asset.occurrences) for asset in raw_assets.values())
    placement_count = sum(len(asset.occurrences) for asset in assets.values())
    _log(
        logger,
        f"Embedded videos found: {len(raw_assets)} file(s), {raw_placement_count} placement(s)",
    )
    if slide_number_map is not None and placement_count != raw_placement_count:
        skipped = raw_placement_count - placement_count
        _log(
            logger,
            f"Hidden-slide mapping active: {placement_count} visible placement(s), {skipped} hidden placement(s) skipped",
        )
    duplicate_descriptions = _duplicate_media_descriptions(assets)
    if duplicate_descriptions:
        preview = "; ".join(duplicate_descriptions[:4])
        suffix = " ..." if len(duplicate_descriptions) > 4 else ""
        _log(logger, f"Repeated video media references detected: {preview}{suffix}")

    source_dir = work_dir / "videos"
    poster_dir = work_dir / "posters"
    _log(logger, "Extracting embedded videos")
    extract_embedded_videos(pptx_path, assets, source_dir)

    sorted_assets = sorted(assets.values(), key=_asset_sort_key)
    total_assets = len(sorted_assets)
    for index, asset in enumerate(sorted_assets, start=1):
        if asset.extracted_path is None:
            continue
        video_name = Path(asset.media_path).name
        raw_poster = poster_dir / f"video_{index:03d}_raw.jpg"
        _log(
            logger,
            f"Video {index}/{total_assets}: extracting poster frame ({video_name})",
        )
        _extract_poster_frame(asset.extracted_path, raw_poster)
        if watermark.enabled:
            with Image.open(raw_poster) as frame_image:
                overlay_path = write_watermark_overlay_image(
                    frame_image.size, watermark
                )
            try:
                prepared_path = source_dir / f"video_{index:03d}_watermarked.mp4"
                _log(
                    logger,
                    f"Video {index}/{total_assets}: transcoding with watermark ({video_name})",
                )

                def log_progress(
                    percent: int, *, current: int = index, total: int = total_assets
                ) -> None:
                    _log(logger, f"Video {current}/{total}: transcoding {percent}%")

                _encode_watermarked_video(
                    asset.extracted_path,
                    overlay_path,
                    prepared_path,
                    encoder_mode=encoder_mode,
                    logger=logger,
                    progress_logger=log_progress,
                )
            finally:
                try:
                    overlay_path.unlink()
                except FileNotFoundError:
                    pass
            poster_path = poster_dir / f"video_{index:03d}_poster.jpg"
            apply_watermark_to_image(raw_poster, poster_path, watermark)
            asset.prepared_path = prepared_path
            asset.poster_path = poster_path
            asset.mime_type = "video/mp4"
            _log(logger, f"Video {index}/{total_assets}: prepared watermarked MP4")
        else:
            asset.prepared_path = asset.extracted_path
            asset.poster_path = raw_poster
            _log(
                logger,
                f"Video {index}/{total_assets}: watermark disabled, keeping source video",
            )
    return assets


def reinsert_videos_into_pptx(
    image_pptx_path: Path,
    assets: dict[str, VideoAsset],
    *,
    logger: Logger | None = None,
) -> int:
    if not assets:
        return 0

    prs = Presentation(str(image_pptx_path))
    slide_count = len(prs.slides)
    inserted = 0
    total_placements = sum(len(asset.occurrences) for asset in assets.values())
    _log(logger, f"Reinserting embedded videos: {total_placements} placement(s)")
    placements = [
        (
            occurrence.slide_number,
            occurrence.occurrence_index,
            asset.media_path,
            asset,
            occurrence,
        )
        for asset in assets.values()
        for occurrence in asset.occurrences
    ]
    for _, _, _, asset, occurrence in sorted(placements):
        if asset.prepared_path is None or asset.poster_path is None:
            continue
        if occurrence.slide_number < 1 or occurrence.slide_number > slide_count:
            _log(
                logger,
                f"Skipping video placement outside rebuilt slide range: slide {occurrence.slide_number}",
            )
            continue
        _log(
            logger,
            f"Video placement {inserted + 1}/{total_placements}: slide {occurrence.slide_number}",
        )
        slide = prs.slides[occurrence.slide_number - 1]
        movie = slide.shapes.add_movie(
            str(asset.prepared_path),
            Emu(occurrence.x),
            Emu(occurrence.y),
            Emu(occurrence.cx),
            Emu(occurrence.cy),
            poster_frame_image=str(asset.poster_path),
            mime_type=asset.mime_type or "video/unknown",
        )
        _apply_original_video_shape_geometry(movie, occurrence)
        inserted += 1

    if not inserted:
        return 0

    _log(logger, "Saving PPTX with reinserted videos")
    fd, temp_name = mkstemp(prefix="pptx_output_watermark_video_", suffix=".pptx")
    os.close(fd)
    temp_output = Path(temp_name)
    try:
        prs.save(str(temp_output))
        shutil.move(str(temp_output), str(image_pptx_path))
    finally:
        try:
            temp_output.unlink()
        except FileNotFoundError:
            pass
    return inserted


def _apply_original_video_shape_geometry(
    movie_shape, occurrence: VideoOccurrence
) -> None:
    element = movie_shape._element
    sp_pr = element.find(qn("p:spPr"))
    xfrm = None if sp_pr is None else sp_pr.find(qn("a:xfrm"))
    if xfrm is not None:
        for attr_name, attr_value in (
            ("rot", occurrence.rotation),
            ("flipH", occurrence.flip_h),
            ("flipV", occurrence.flip_v),
        ):
            if attr_value:
                xfrm.set(attr_name, attr_value)

    if not occurrence.crop:
        return
    blip_fill = element.find(qn("p:blipFill"))
    if blip_fill is None:
        return
    src_rect = blip_fill.find(qn("a:srcRect"))
    if src_rect is None:
        src_rect = OxmlElement("a:srcRect")
        blip = blip_fill.find(qn("a:blip"))
        if blip is not None:
            blip.addnext(src_rect)
        else:
            blip_fill.insert(0, src_rect)
    for key in ("l", "t", "r", "b"):
        if key in occurrence.crop:
            src_rect.set(key, occurrence.crop[key])


def export_sidecar_videos(
    output_pptx_path: Path,
    assets: dict[str, VideoAsset],
) -> VideoSidecarInfo | None:
    if not assets:
        return None
    sidecar_dir = _available_sidecar_dir(output_pptx_path)
    sidecar_dir.mkdir(parents=True, exist_ok=True)
    manifest: list[dict[str, object]] = []
    for asset in sorted(assets.values(), key=lambda item: item.media_path):
        if asset.prepared_path is None:
            continue
        stem = (
            re.sub(r"[^A-Za-z0-9._-]+", "_", Path(asset.media_path).stem).strip("_")
            or "video"
        )
        suffix = asset.prepared_path.suffix or ".mp4"
        for occurrence in asset.occurrences:
            filename = f"S{occurrence.slide_number:03d}_{occurrence.occurrence_index:02d}_{stem}{suffix}"
            target_path = sidecar_dir / filename
            shutil.copy2(asset.prepared_path, target_path)
            manifest.append(
                {
                    "slide_number": occurrence.slide_number,
                    "occurrence_index": occurrence.occurrence_index,
                    "shape_name": occurrence.shape_name,
                    "source_media_path": asset.media_path,
                    "output_file": target_path.name,
                    "x": occurrence.x,
                    "y": occurrence.y,
                    "cx": occurrence.cx,
                    "cy": occurrence.cy,
                }
            )
    manifest_path = sidecar_dir / "manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "pptx_output": str(output_pptx_path),
                "videos": manifest,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return VideoSidecarInfo(directory=sidecar_dir, manifest_path=manifest_path)


def _available_sidecar_dir(output_pptx_path: Path) -> Path:
    base = output_pptx_path.with_name(f"{output_pptx_path.stem}_videos")
    if not base.exists():
        return base
    for index in range(2, 1000):
        candidate = output_pptx_path.with_name(
            f"{output_pptx_path.stem}_videos_{index}"
        )
        if not candidate.exists():
            return candidate
    raise RuntimeError(
        f"Unable to find available sidecar directory beside {output_pptx_path}"
    )


def _audio_rel_ids(pic: ET.Element) -> list[str]:
    """Return relationship ids for embedded/linked audio in a p:pic shape."""
    nv_pic = pic.find("./p:nvPicPr", NS)
    if nv_pic is None:
        return []
    nv_pr = nv_pic.find("./p:nvPr", NS)
    if nv_pr is None:
        return []
    rel_ids: list[str] = []
    audio_file = nv_pr.find("a:audioFile", NS)
    if audio_file is not None:
        rel_id = audio_file.attrib.get(f"{R_NS}link")
        if rel_id:
            rel_ids.append(rel_id)
    media_ref = nv_pr.find(".//p14:media", NS)
    if media_ref is not None:
        rel_id = media_ref.attrib.get(f"{R_NS}embed")
        if rel_id:
            rel_ids.append(rel_id)
    return rel_ids


def _parse_relationships_with_external(
    zf: ZipFile, slide_path: str
) -> tuple[dict[str, str], dict[str, str]]:
    """Return (internal_targets, external_targets) for a slide's rels."""
    rels_path = _relationships_path_for(slide_path)
    if rels_path not in zf.namelist():
        return {}, {}
    root = ET.fromstring(zf.read(rels_path))
    internal: dict[str, str] = {}
    external: dict[str, str] = {}
    for rel in root.findall("pr:Relationship", REL_NS):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if not rel_id or not target:
            continue
        if rel.attrib.get("TargetMode") == "External":
            external[rel_id] = target
        else:
            internal[rel_id] = _resolve_zip_target(slide_path, target)
    return internal, external


def scan_fidelity_warnings(pptx_path: Path) -> list[str]:
    """Detect media content silently lost by the image-PPTX rebuild path.

    Returns a list of human-readable (Chinese) warning strings. Each describes
    a category of content that the image-PPTX + video-reinsert pipeline does
    not preserve, so callers can surface them to the user instead of silently
    dropping the content.

    Detected categories:
    - Embedded/linked audio shapes (a:audioFile / p14:media on audio)
    - Linked (external) video references
    - Videos inside grouped shapes (coordinates are group-local)
    - Video files with extensions outside VIDEO_EXTENSIONS
    - Videos on slide masters/layouts (not scanned by the rebuild path)
    """
    warnings: list[str] = []
    audio_count = 0
    linked_video_count = 0
    grouped_video_count = 0
    unsupported_video_count = 0

    with ZipFile(pptx_path) as zf:
        names = set(zf.namelist())

        # --- slide masters/layouts videos (not scanned by rebuild) ---
        master_video_count = 0
        for member in names:
            if not (
                member.startswith("ppt/slideMasters/slideMaster")
                or member.startswith("ppt/slideLayouts/slideLayout")
            ) or not member.endswith(".xml"):
                continue
            try:
                root = ET.fromstring(zf.read(member))
            except ET.ParseError:
                continue
            for pic in root.findall(".//p:pic", NS):
                if _video_rel_ids(pic) or _audio_rel_ids(pic):
                    master_video_count += 1
        if master_video_count:
            warnings.append(
                f"母版/版式上有 {master_video_count} 个媒体对象不会出现在导出页面中。"
            )

        # --- slide-level scan ---
        slide_paths = _presentation_ordered_slide_paths(zf)
        for slide_path in slide_paths:
            internal, external = _parse_relationships_with_external(zf, slide_path)
            try:
                root = ET.fromstring(zf.read(slide_path))
            except ET.ParseError:
                continue

            # detect grouping by finding p:pic inside p:grpSp
            for grp in root.findall(".//p:grpSp", NS):
                for pic in grp.findall(".//p:pic", NS):
                    if _video_rel_ids(pic):
                        grouped_video_count += 1

            for pic in root.findall(".//p:pic", NS):
                # audio shapes
                if _audio_rel_ids(pic):
                    audio_count += 1
                    continue

                video_rel_ids = _video_rel_ids(pic)
                if not video_rel_ids:
                    continue

                for rel_id in video_rel_ids:
                    if rel_id in external:
                        linked_video_count += 1
                        continue
                    target = internal.get(rel_id, "")
                    if target and Path(target).suffix.lower() not in VIDEO_EXTENSIONS:
                        unsupported_video_count += 1

    if audio_count:
        warnings.append(
            f"检测到 {audio_count} 个嵌入/链接音频（背景音乐、配音等），"
            "图片 PPTX 模式不保留音频，导出后将静默丢失。"
        )
    if linked_video_count:
        warnings.append(
            f"检测到 {linked_video_count} 个外部链接视频，"
            "图片 PPTX 模式仅保留嵌入视频，链接视频将丢失。"
        )
    if grouped_video_count:
        warnings.append(
            f"检测到 {grouped_video_count} 个组合（group）内的视频，"
            "其坐标为组内局部坐标系，回填后位置可能错乱。"
        )
    if unsupported_video_count:
        warnings.append(
            f"检测到 {unsupported_video_count} 个非白名单容器的视频文件，"
            "将不被识别和回填。"
        )
    return warnings
