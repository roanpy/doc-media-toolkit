from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from pptx_output_watermark.models import WatermarkOptions
from pptx_output_watermark.pptx_video_support import (
    replace_video_media_in_pptx,
    scan_embedded_videos,
    watermark_videos_in_editable_pptx,
)


def _make_poster(path: Path) -> Path:
    Image.new("RGB", (16, 16), (10, 20, 30)).save(path)
    return path


def _build_pptx_with_video(pptx_path: Path) -> str:
    """Build a minimal PPTX with one embedded video; return its media path."""
    video_path = pptx_path.with_suffix(".mp4")
    video_path.write_bytes(b"synthetic-test-video")
    poster = pptx_path.with_suffix(".poster.png")
    _make_poster(poster)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_movie(
        str(video_path),
        Emu(1000000),
        Emu(1000000),
        Emu(2000000),
        Emu(2000000),
        poster_frame_image=str(poster),
        mime_type="video/mp4",
    )
    prs.save(str(pptx_path))
    # find the media path
    assets = scan_embedded_videos(pptx_path)
    assert assets, "expected at least one video"
    return next(iter(assets))


class ReplaceVideoMediaInPptxTest(unittest.TestCase):
    def test_replaces_media_bytes_while_preserving_slide_xml(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            d_path = Path(d)
            pptx_path = d_path / "input.pptx"
            media_path = _build_pptx_with_video(pptx_path)

            # snapshot original slide XML and media bytes
            with ZipFile(pptx_path) as zf:
                original_slide1 = zf.read("ppt/slides/slide1.xml")
                original_media = zf.read(media_path)

            # craft a fake "watermarked" replacement that differs from original
            replacement = d_path / "fake_watermarked.mp4"
            replacement.write_bytes(original_media + b"\x00WATERMARK\x00")
            self.assertNotEqual(replacement.read_bytes(), original_media)

            output = d_path / "output.pptx"
            replace_video_media_in_pptx(
                pptx_path, {media_path: replacement}, output_path=output
            )

            with ZipFile(output) as zf:
                # slide XML byte-identical => timing/playback/settings preserved
                self.assertEqual(zf.read("ppt/slides/slide1.xml"), original_slide1)
                # media bytes replaced
                self.assertEqual(zf.read(media_path), replacement.read_bytes())

    def test_replacement_defaults_to_overwrite_in_place(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            d_path = Path(d)
            pptx_path = d_path / "input.pptx"
            media_path = _build_pptx_with_video(pptx_path)
            with ZipFile(pptx_path) as zf:
                original_slide1 = zf.read("ppt/slides/slide1.xml")

            replacement = d_path / "wm.mp4"
            replacement.write_bytes(b"FAKE_MEDIA_BYTES")
            result = replace_video_media_in_pptx(pptx_path, {media_path: replacement})

            self.assertEqual(result, pptx_path)
            with ZipFile(pptx_path) as zf:
                self.assertEqual(zf.read(media_path), b"FAKE_MEDIA_BYTES")
                self.assertEqual(zf.read("ppt/slides/slide1.xml"), original_slide1)

    def test_no_replacements_copies_all_entries_unchanged(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            d_path = Path(d)
            pptx_path = d_path / "input.pptx"
            _build_pptx_with_video(pptx_path)
            with ZipFile(pptx_path) as zf:
                original_entries = {n: zf.read(n) for n in zf.namelist()}

            output = d_path / "output.pptx"
            replace_video_media_in_pptx(pptx_path, {}, output_path=output)

            with ZipFile(output) as zf:
                for name, data in original_entries.items():
                    self.assertEqual(zf.read(name), data, f"entry changed: {name}")


class WatermarkVideosInEditablePptxTest(unittest.TestCase):
    """End-to-end (with ffmpeg mocked) of the editable-PPTX video watermark
    path: scans, 'encodes', replaces media, and preserves slide XML."""

    def test_watermarks_videos_preserving_slide_xml(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            d_path = Path(d)
            pptx_path = d_path / "input.pptx"
            media_path = _build_pptx_with_video(pptx_path)

            with ZipFile(pptx_path) as zf:
                original_slide1 = zf.read("ppt/slides/slide1.xml")
                original_timing_count = original_slide1.count(b"<p:timing")

            # fake encoder: write a distinct byte pattern so we can detect replacement
            def fake_encode(source_path, output_path, watermark, **kwargs):  # noqa: ANN001
                output_path.parent.mkdir(parents=True, exist_ok=True)
                output_path.write_bytes(b"WATERMARKED_VIDEO_BYTES")
                return output_path

            output = d_path / "output.pptx"
            with patch(
                "pptx_output_watermark.pptx_video_support.watermark_video_file",
                side_effect=fake_encode,
            ):
                watermark_videos_in_editable_pptx(
                    pptx_path,
                    output,
                    WatermarkOptions(enabled=True),
                )

            with ZipFile(output) as zf:
                # timing tree preserved (add_movie always creates one <p:timing>)
                self.assertEqual(
                    zf.read("ppt/slides/slide1.xml").count(b"<p:timing"),
                    original_timing_count,
                )
                self.assertEqual(zf.read("ppt/slides/slide1.xml"), original_slide1)
                # media bytes replaced with the fake watermarked content
                self.assertEqual(zf.read(media_path), b"WATERMARKED_VIDEO_BYTES")

    def test_no_videos_copies_pptx_unchanged(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            d_path = Path(d)
            pptx_path = d_path / "input.pptx"
            # build a PPTX with no video
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(str(pptx_path))
            with ZipFile(pptx_path) as zf:
                original_entries = {n: zf.read(n) for n in zf.namelist()}

            output = d_path / "output.pptx"
            watermark_videos_in_editable_pptx(
                pptx_path,
                output,
                WatermarkOptions(enabled=True),
            )
            with ZipFile(output) as zf:
                self.assertEqual(set(zf.namelist()), set(original_entries))
                for name, data in original_entries.items():
                    self.assertEqual(zf.read(name), data)


if __name__ == "__main__":
    unittest.main()
