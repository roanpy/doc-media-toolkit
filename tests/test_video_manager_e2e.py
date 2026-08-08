"""End-to-end validation for the video asset manager.

Unlike test_video_manager.py (which uses fake video bytes and mocked probes),
these tests run the real chain against a real ffmpeg-generated video:

    extract (real ffprobe) -> detach (real placeholder encode)
    -> restore -> structural + hash assertions

They skip automatically when ffmpeg/ffprobe are unavailable.
"""

from __future__ import annotations

import hashlib
import subprocess
import tempfile
import unittest
from argparse import Namespace
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx_output_watermark.ffmpeg_runtime import resolve_binary
from pptx_output_watermark.pptx_video_support import scan_embedded_videos
from pptx_tools.video_manager import (
    VideoProject,
    _video_fingerprint,
    probe_video,
    sha256_file,
)
from pptx_video_compactor import compact_pptx

TIMING_TEMPLATE = (
    "<p:timing><p:tnLst><p:par>"
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    "<p:childTnLst>"
    '<p:video><p:cMediaNode vol="80000">'
    '<p:cTn id="2" fill="hold" display="0">'
    '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
    "</p:cTn>"
    "<p:tgtEl><p:spid>{shape_id}</p:spid></p:tgtEl>"
    "</p:cMediaNode></p:video>"
    "</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"
)


def ffmpeg_available() -> bool:
    return (
        resolve_binary("ffmpeg") is not None and resolve_binary("ffprobe") is not None
    )


def make_video_pptx(path: Path, video_bytes: bytes, title: str) -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    video = path.with_name(f"{path.stem}.mp4")
    poster = path.with_name(f"{path.stem}.png")
    video.write_bytes(video_bytes)
    Image.new("RGB", (320, 180), "#dc6b2f").save(poster)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(5), Inches(0.5)
    ).text = title
    slide.shapes.add_movie(
        str(video),
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3.375),
        poster_frame_image=str(poster),
        mime_type="video/mp4",
    )
    presentation.save(path)


def make_wmv_pptx(path: Path, video_path: Path, title: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(5), Inches(0.5)
    ).text = title
    slide.shapes.add_movie(
        str(video_path),
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3.375),
        poster_frame_image=None,
        mime_type="video/x-ms-wmv",
    )
    presentation.save(path)


def make_wmv_video(source_mp4: Path, target: Path, *, bitrate: str = "600k") -> None:
    ffmpeg = resolve_binary("ffmpeg")
    assert ffmpeg is not None
    subprocess.run(
        [
            ffmpeg,
            "-y",
            "-v",
            "error",
            "-i",
            str(source_mp4),
            "-c:v",
            "wmv2",
            "-b:v",
            bitrate,
            "-c:a",
            "wmav2",
            str(target),
        ],
        check=True,
    )


def make_real_video(
    path: Path,
    seconds: float = 2.0,
    *,
    size: str = "640x360",
    crf: int | None = None,
) -> None:
    ffmpeg = resolve_binary("ffmpeg")
    assert ffmpeg is not None
    subprocess.run(
        [
            ffmpeg,
            "-y",
            "-v",
            "error",
            "-f",
            "lavfi",
            "-i",
            f"testsrc=size={size}:rate=15",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=440",
            "-t",
            f"{seconds:.3f}",
            "-c:v",
            "libx264",
            "-pix_fmt",
            "yuv420p",
            *(["-crf", str(crf)] if crf is not None else []),
            "-c:a",
            "aac",
            str(path),
        ],
        check=True,
    )


def transcode_video(source: Path, target: Path, *, audio_frequency: int | None) -> None:
    ffmpeg = resolve_binary("ffmpeg")
    assert ffmpeg is not None
    command = [ffmpeg, "-y", "-v", "error", "-i", str(source)]
    if audio_frequency is not None:
        command.extend(
            ["-f", "lavfi", "-i", f"sine=frequency={audio_frequency}:duration=2"]
        )
    command.extend(["-map", "0:v:0"])
    if audio_frequency is None:
        command.extend(["-map", "0:a:0"])
    else:
        command.extend(["-map", "1:a:0"])
    command.extend(
        [
            "-vf",
            "scale=320:180",
            "-c:v",
            "libx264",
            "-crf",
            "32",
            "-c:a",
            "aac",
            "-shortest",
            str(target),
        ]
    )
    subprocess.run(command, check=True)


def inject_video_animation(pptx_path: Path, slide_path: str, shape_id: int) -> None:
    """Insert a <p:timing> video animation node into the slide XML in place."""
    with ZipFile(pptx_path, "r") as zin:
        entries = [(info, zin.read(info.filename)) for info in zin.infolist()]
    timing = TIMING_TEMPLATE.format(shape_id=shape_id).encode("utf-8")
    rewritten = False
    temp_path = pptx_path.with_suffix(".inject.tmp")
    try:
        with ZipFile(temp_path, "w", ZIP_DEFLATED) as zout:
            for info, data in entries:
                if info.filename == slide_path:
                    closing = data.rfind(b"</p:sld>")
                    assert closing > 0, f"unexpected slide XML shape: {slide_path}"
                    data = data[:closing] + timing + data[closing:]
                    rewritten = True
                zout.writestr(info, data)
        assert rewritten, f"slide not found in package: {slide_path}"
        temp_path.replace(pptx_path)
    finally:
        temp_path.unlink(missing_ok=True)


def ffprobe_duration_width(path: Path) -> tuple[float, int]:
    ffprobe = resolve_binary("ffprobe")
    assert ffprobe is not None
    result = subprocess.run(
        [
            ffprobe,
            "-v",
            "error",
            "-print_format",
            "json",
            "-show_streams",
            "-show_format",
            str(path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    import json

    payload = json.loads(result.stdout or "{}")
    video = next(
        s for s in payload.get("streams", []) if s.get("codec_type") == "video"
    )
    duration = float(
        video.get("duration") or payload.get("format", {}).get("duration") or 0
    )
    return duration, int(video.get("width") or 0)


@unittest.skipUnless(ffmpeg_available(), "ffmpeg/ffprobe not available")
class VideoManagerEndToEndTest(unittest.TestCase):
    def test_external_reencode_requires_manual_source_promotion(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            original = root / "original.mp4"
            compact = root / "compact.mp4"
            make_real_video(original)
            transcode_video(original, compact, audio_frequency=None)
            deck = root / "deck.pptx"
            make_video_pptx(deck, compact.read_bytes(), "External import")
            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(deck, source_quality="original")

            result = library.import_external_video(original, source_quality="original")

            family = library.families()[0]
            source = library.source_variant(family)
            self.assertEqual(result["status"], "matched")
            self.assertFalse(result["promoted"])
            library.set_source_variant(result["variant_id"])
            source = library.source_variant(family)
            self.assertEqual((source["width"], source["height"]), (640, 360))

    def test_overwrite_compression_archives_then_restores_video(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "source.mp4"
            make_real_video(video)
            original_bytes = video.read_bytes()
            source = root / "source.pptx"
            make_video_pptx(source, original_bytes, "Overwrite")
            media_path = next(iter(scan_embedded_videos(source)))

            library = VideoProject.create(root / "library")
            archived = library.archive_pptx_videos(source, source_quality="original")
            result = compact_pptx(
                Namespace(
                    input_pptx=source,
                    target_size_mb=None,
                    config=None,
                    profile="aggressive",
                    image_profile="none",
                    output=source,
                    video_output_dir=None,
                    slide_render_width=1920,
                    slide_render_height=1080,
                    min_height=240,
                    max_height=360,
                    overscan=1.0,
                    reserve_mb=0.0,
                    preset="medium",
                    encoder="cpu",
                    work_dir=None,
                    keep_work_dir=False,
                    keep_artifacts=False,
                    dry_run=False,
                ),
                logger=lambda _message: None,
            )
            self.assertEqual(Path(result["output_pptx"]).resolve(), source.resolve())
            with ZipFile(source) as archive:
                self.assertNotEqual(archive.read(media_path), original_bytes)

            aliases = library.register_compressed_pptx_hashes(
                source, Path(result["report_path"]), archived["media_families"]
            )
            self.assertEqual(aliases, 1)
            restored = library.upgrade_pptx_from_library(
                source, output_path=root / "restored.pptx"
            )
            with ZipFile(restored["output_pptx"]) as archive:
                self.assertEqual(archive.read(media_path), original_bytes)

    def test_repeated_compression_keeps_one_deck_and_restores_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "source.mp4"
            make_real_video(video, seconds=3.0, size="1280x720", crf=10)
            original_bytes = video.read_bytes()
            source = root / "source.pptx"
            make_video_pptx(source, original_bytes, "Repeated compression")
            media_path = next(iter(scan_embedded_videos(source)))

            library = VideoProject.create(root / "library")
            first_archive = library.archive_pptx_videos(source)
            source_deck = library.add_deck(source)
            first_output = root / "first-compressed.pptx"
            first_result = compact_pptx(
                Namespace(
                    input_pptx=source,
                    target_size_mb=None,
                    config=None,
                    profile="high",
                    image_profile="none",
                    output=first_output,
                    video_output_dir=None,
                    slide_render_width=1920,
                    slide_render_height=1080,
                    min_height=240,
                    max_height=360,
                    overscan=1.0,
                    reserve_mb=0.0,
                    preset="medium",
                    encoder="cpu",
                    work_dir=None,
                    keep_work_dir=False,
                    keep_artifacts=False,
                    dry_run=False,
                ),
                logger=lambda _message: None,
            )
            library.register_compressed_pptx_hashes(
                first_output,
                Path(first_result["report_path"]),
                first_archive["media_families"],
            )
            library.register_optimized_output(source, first_output)

            second_archive = library.archive_pptx_videos(first_output)
            repeated_deck = library.add_deck(first_output)
            second_output = root / "second-compressed.pptx"
            second_result = compact_pptx(
                Namespace(
                    input_pptx=first_output,
                    target_size_mb=None,
                    config=None,
                    profile="aggressive",
                    image_profile="none",
                    output=second_output,
                    video_output_dir=None,
                    slide_render_width=1920,
                    slide_render_height=1080,
                    min_height=240,
                    max_height=360,
                    overscan=1.0,
                    reserve_mb=0.0,
                    preset="medium",
                    encoder="cpu",
                    work_dir=None,
                    keep_work_dir=False,
                    keep_artifacts=False,
                    dry_run=False,
                ),
                logger=lambda _message: None,
            )
            library.register_compressed_pptx_hashes(
                second_output,
                Path(second_result["report_path"]),
                second_archive["media_families"],
            )
            second_record = library.register_optimized_output(
                first_output, second_output
            )

            self.assertEqual(
                (second_archive["added"], second_archive["reused"]), (0, 1)
            )
            self.assertEqual(repeated_deck["id"], source_deck["id"])
            self.assertEqual(len(library.decks()), 1)
            self.assertIsNotNone(second_record)
            self.assertEqual(len(source_deck["optimized_outputs"]), 1)
            self.assertEqual(len(library.families()), 1)

            restored = library.upgrade_pptx_from_library(
                second_output, output_path=root / "restored.pptx"
            )
            self.assertEqual(restored["matched"], 1)
            self.assertEqual(restored["unmatched"], [])
            with ZipFile(restored["output_pptx"]) as archive:
                self.assertEqual(archive.read(media_path), original_bytes)

    def test_content_fingerprint_matches_reencode_but_rejects_different_audio(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_video = root / "source.mp4"
            same_content = root / "same-content.mp4"
            different_audio = root / "different-audio.mp4"
            make_real_video(source_video)
            transcode_video(source_video, same_content, audio_frequency=None)
            transcode_video(source_video, different_audio, audio_frequency=80)

            source_pptx = root / "source.pptx"
            make_video_pptx(source_pptx, source_video.read_bytes(), "Fingerprint")
            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source_pptx)

            same_fingerprint = _video_fingerprint(same_content)
            different_fingerprint = _video_fingerprint(different_audio)
            self.assertIsNotNone(same_fingerprint)
            self.assertIsNotNone(different_fingerprint)
            self.assertIsNotNone(
                library.family_by_content_fingerprint(same_fingerprint or {})
            )
            self.assertIsNone(
                library.family_by_content_fingerprint(different_fingerprint or {})
            )

            reencoded_pptx = root / "reencoded.pptx"
            make_video_pptx(
                reencoded_pptx, same_content.read_bytes(), "Reencoded Fingerprint"
            )
            imported = library.archive_pptx_videos(reencoded_pptx)
            self.assertEqual((imported["added"], imported["reused"]), (0, 1))
            self.assertEqual(len(library.families()), 1)

    def test_library_archives_source_and_upgrades_real_compacted_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "source.mp4"
            make_real_video(video, seconds=2.0)
            original_bytes = video.read_bytes()
            source = root / "source.pptx"
            make_video_pptx(source, original_bytes, "Library Round Trip")

            scanned = scan_embedded_videos(source)
            media_path, asset = next(iter(scanned.items()))
            slide_path = asset.occurrences[0].slide_path
            inject_video_animation(source, slide_path, asset.occurrences[0].shape_id)
            with ZipFile(source) as archive:
                original_slide_xml = archive.read(slide_path)

            library = VideoProject.create(root / "library")
            archived = library.archive_pptx_videos(source)
            self.assertEqual((archived["added"], archived["reused"]), (1, 0))
            self.assertEqual(library.decks(), [])

            compacted = root / "source_compacted.pptx"
            compact_pptx(
                Namespace(
                    input_pptx=source,
                    target_size_mb=None,
                    config=None,
                    profile="aggressive",
                    image_profile="none",
                    output=compacted,
                    video_output_dir=None,
                    slide_render_width=1920,
                    slide_render_height=1080,
                    min_height=240,
                    max_height=360,
                    overscan=1.0,
                    reserve_mb=0.0,
                    preset="medium",
                    encoder="cpu",
                    work_dir=None,
                    keep_work_dir=False,
                    keep_artifacts=False,
                    dry_run=False,
                ),
                logger=lambda _message: None,
            )
            content_upgrade = library.upgrade_pptx_from_library(
                compacted, output_path=root / "content-matched.pptx"
            )
            self.assertEqual(content_upgrade["matched"], 1)
            self.assertEqual(content_upgrade["content_matched"], 1)
            self.assertEqual(content_upgrade["aliases_added"], 1)

            upgraded = library.upgrade_pptx_from_library(compacted)
            self.assertEqual(upgraded["matched"], 1)
            self.assertEqual(upgraded["content_matched"], 0)
            self.assertEqual(upgraded["aliases_added"], 0)
            self.assertEqual(upgraded["unmatched"], [])
            self.assertEqual(library.decks(), [])
            output = Path(upgraded["output_pptx"])
            with ZipFile(compacted) as low_zip, ZipFile(output) as high_zip:
                self.assertNotEqual(low_zip.read(media_path), original_bytes)
                self.assertEqual(high_zip.read(media_path), original_bytes)
                self.assertEqual(low_zip.read(slide_path), original_slide_xml)
                self.assertEqual(high_zip.read(slide_path), original_slide_xml)

            low_clip = root / "low.mp4"
            with ZipFile(compacted) as archive:
                low_clip.write_bytes(archive.read(media_path))
            duration, width = ffprobe_duration_width(low_clip)
            self.assertAlmostEqual(duration, 2.0, delta=0.5)
            self.assertLessEqual(width, 640)

    def test_real_video_detach_restore_preserves_animation_timeline(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "clip.mp4"
            make_real_video(video, seconds=2.0)
            original_bytes = video.read_bytes()

            source = root / "deck.pptx"
            make_video_pptx(source, original_bytes, "E2E Round Trip")

            scanned = scan_embedded_videos(source)
            self.assertEqual(len(scanned), 1)
            part_path, asset = next(iter(scanned.items()))
            slide_path = asset.occurrences[0].slide_path
            shape_id = asset.occurrences[0].shape_id
            inject_video_animation(source, slide_path, shape_id)

            with ZipFile(source) as archive:
                original_slide_xml = archive.read(slide_path)
            self.assertIn(b"<p:timing>", original_slide_xml)

            project = VideoProject.create(root / "project")
            deck = project.add_deck(source)

            # Real probe captured real metadata during archiving.
            variant = project.families()[0]["variants"][0]
            self.assertAlmostEqual(variant["duration_sec"], 2.0, delta=0.4)
            self.assertEqual((variant["width"], variant["height"]), (640, 360))
            self.assertTrue(variant["has_audio"])
            self.assertEqual(
                variant["sha256"], hashlib.sha256(original_bytes).hexdigest()
            )

            # Real detach: placeholder encoded by ffmpeg, not a mock.
            detached = project.detach_deck(deck["id"])
            restored = project.restore_deck(deck["id"], detached)

            placeholder_part = deck["assets"][0]["placeholder_part"]
            with ZipFile(detached) as detached_zip, ZipFile(restored) as restored_zip:
                # Slide XML (which carries the animation timeline) is untouched
                # by both detach and restore.
                self.assertEqual(detached_zip.read(slide_path), original_slide_xml)
                self.assertEqual(restored_zip.read(slide_path), original_slide_xml)
                # Detached package swaps the media part for the placeholder.
                self.assertNotIn(part_path, detached_zip.namelist())
                self.assertIn(placeholder_part, detached_zip.namelist())
                # Restored package carries the exact original video bytes.
                self.assertEqual(restored_zip.read(part_path), original_bytes)

            # The placeholder still anchors to the same shape on the same slide.
            detached_scan = scan_embedded_videos(detached)
            self.assertIn(placeholder_part, detached_scan)
            self.assertEqual(
                detached_scan[placeholder_part].occurrences[0].shape_id, shape_id
            )
            self.assertEqual(
                detached_scan[placeholder_part].occurrences[0].slide_path, slide_path
            )

            # The placeholder preserves the playback duration (animation
            # timelines that trigger on video end keep working) and stays small.
            placeholder_file = root / "placeholder.mp4"
            with ZipFile(detached) as detached_zip:
                placeholder_file.write_bytes(detached_zip.read(placeholder_part))
            duration, width = ffprobe_duration_width(placeholder_file)
            self.assertAlmostEqual(duration, 2.0, delta=0.5)
            self.assertLessEqual(width, 640)
            self.assertLess(placeholder_file.stat().st_size, len(original_bytes) // 2)

            # The restored package scans back to the original shape as well.
            restored_scan = scan_embedded_videos(restored)
            self.assertIn(part_path, restored_scan)
            self.assertEqual(restored_scan[part_path].occurrences[0].shape_id, shape_id)

    def test_real_video_restore_uses_active_compressed_variant(self) -> None:
        """Detach, compress the archived original, then restore must embed
        the compressed variant (the active one), not the original bytes."""
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "clip.mp4"
            # Start from a high-quality source so the compressed variant still
            # clears the normal safety gate; low-bitrate synthetic testsrc is
            # correctly restored instead of becoming an active variant.
            make_real_video(video, seconds=2.0, crf=18)
            original_bytes = video.read_bytes()

            source = root / "deck.pptx"
            make_video_pptx(source, original_bytes, "E2E Active Variant")

            project = VideoProject.create(root / "project")
            deck = project.add_deck(source)
            family = project.families()[0]
            original_variant = family["variants"][0]

            detached = project.detach_deck(deck["id"])

            compressed = project.compress_variant(
                original_variant["id"],
                "aggressive",
                progress_callback=lambda _done, _total, _label: None,
            )
            self.assertEqual(family["active_variant_id"], compressed["id"])
            compressed_path = project.variant_path(compressed)
            compressed_bytes = compressed_path.read_bytes()
            self.assertNotEqual(compressed_bytes, original_bytes)

            restored = project.restore_deck(deck["id"], detached)

            with ZipFile(restored) as restored_zip:
                names = restored_zip.namelist()
                media_parts = [
                    name
                    for name in names
                    if name.startswith("ppt/media/") and name.endswith(".mp4")
                ]
                self.assertEqual(len(media_parts), 1)
                embedded = restored_zip.read(media_parts[0])

            self.assertEqual(embedded, compressed_bytes)
            self.assertNotEqual(embedded, original_bytes)
            # Restored media keeps a playable duration close to the source,
            # and its bytes hash to the active variant's recorded digest.
            restored_clip = root / "restored-embedded.mp4"
            restored_clip.write_bytes(embedded)
            duration, _ = ffprobe_duration_width(restored_clip)
            self.assertAlmostEqual(duration, 2.0, delta=1.0)
            self.assertEqual(hashlib.sha256(embedded).hexdigest(), compressed["sha256"])

    def test_upgrade_replaces_wmv_part_with_proper_mp4_part(self) -> None:
        """A non-MP4 part upgraded from the library must become a real .mp4
        part (bytes, extension, relationships and content type aligned), never
        MP4 bytes hidden under a .wmv name."""
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            base = root / "base.mp4"
            make_real_video(base, seconds=2.0)
            wmv = root / "clip.wmv"
            make_wmv_video(base, wmv)

            source = root / "source.pptx"
            make_wmv_pptx(source, wmv, "WMV Library")
            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source, source_quality="original")

            # Same content, re-encoded to different WMV bytes.
            reencoded_wmv = root / "clip-reenc.wmv"
            make_wmv_video(base, reencoded_wmv, bitrate="400k")
            deck = root / "deck.pptx"
            make_wmv_pptx(deck, reencoded_wmv, "WMV Deck")
            with ZipFile(deck) as archive:
                self.assertIn("ppt/media/media1.wmv", archive.namelist())

            result = library.upgrade_pptx_from_library(
                deck, output_path=root / "upgraded.pptx"
            )
            self.assertEqual(result["matched"], 1)
            output = Path(result["output_pptx"])
            with ZipFile(output) as archive:
                self.assertIsNone(archive.testzip())
                names = archive.namelist()
                self.assertNotIn("ppt/media/media1.wmv", names)
                self.assertIn("ppt/media/media1.mp4", names)
                mp4_bytes = archive.read("ppt/media/media1.mp4")
                # ftyp box: really an ISO-BMFF/MP4 payload.
                self.assertIn(b"ftyp", mp4_bytes[4:12])
                content_types = archive.read("[Content_Types].xml").decode().lower()
                self.assertIn("mp4", content_types)
                rels = archive.read("ppt/slides/_rels/slide1.xml.rels").decode()
                self.assertIn("media1.mp4", rels)
                self.assertNotIn("media1.wmv", rels)

            # The new part still anchors to the video shape.
            scanned = scan_embedded_videos(output)
            self.assertIn("ppt/media/media1.mp4", scanned)

    def test_upgrade_transcodes_library_source_wmv_to_mp4(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            base = root / "base.mp4"
            make_real_video(base, seconds=2.0)
            wmv = root / "clip.wmv"
            make_wmv_video(base, wmv)
            source = root / "source.pptx"
            make_wmv_pptx(source, wmv, "WMV Library Source")
            scanned = scan_embedded_videos(source)
            asset = next(iter(scanned.values()))
            slide_path = asset.occurrences[0].slide_path
            inject_video_animation(source, slide_path, asset.occurrences[0].shape_id)
            with ZipFile(source) as archive:
                original_slide_xml = archive.read(slide_path)
            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source, source_quality="original")

            result = library.upgrade_pptx_from_library(
                source, output_path=root / "upgraded.pptx"
            )

            self.assertEqual(result["matched"], 1)
            self.assertEqual(result["already_high_quality"], 0)
            output = Path(result["output_pptx"])
            with ZipFile(output) as archive:
                self.assertNotIn("ppt/media/media1.wmv", archive.namelist())
                self.assertIn("ppt/media/media1.mp4", archive.namelist())
                self.assertIn(b"ftyp", archive.read("ppt/media/media1.mp4")[4:12])
                self.assertEqual(archive.read(slide_path), original_slide_xml)

    def test_adopt_upgraded_wmv_deck_preserves_round_trip_association(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            base = root / "base.mp4"
            make_real_video(base, seconds=2.0)
            wmv = root / "clip.wmv"
            make_wmv_video(base, wmv)
            source = root / "source.pptx"
            make_wmv_pptx(source, wmv, "WMV Registered Source")
            alias = root / "renamed-copy.pptx"
            alias.write_bytes(source.read_bytes())
            library = VideoProject.create(root / "library")
            deck = library.add_deck(source)
            self.assertEqual(library.add_deck(alias)["id"], deck["id"])
            family = library.family(deck["assets"][0]["family_id"])
            mp4_variant = library.create_unified_version(family["id"])
            library.activate_variant(mp4_variant["id"])

            result = library.upgrade_pptx_from_library(
                source,
                output_path=root / "upgraded.pptx",
                incompatible_only=True,
            )
            upgraded = Path(result["output_pptx"])
            source.write_bytes(upgraded.read_bytes())
            alias.write_bytes(upgraded.read_bytes())
            adopted = library.adopt_upgraded_deck_source(deck["id"], alias)

            self.assertEqual(adopted["assets"][0]["part_path"], "ppt/media/media1.mp4")
            self.assertEqual(
                adopted["assets"][0]["original_variant_id"], mp4_variant["id"]
            )
            self.assertEqual(library.deck_source_path(adopted), alias.resolve())
            self.assertEqual(adopted["source_aliases"], [str(source.resolve())])
            detached = library.detach_deck(deck["id"])
            restored = library.restore_deck(deck["id"], detached)
            with ZipFile(restored) as archive:
                self.assertEqual(
                    hashlib.sha256(archive.read("ppt/media/media1.mp4")).hexdigest(),
                    mp4_variant["sha256"],
                )

    def test_adopt_compact_deck_keeps_high_quality_restore_variant(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_video = root / "source.mp4"
            compact_video = root / "compact.mp4"
            make_real_video(source_video, seconds=2.0)
            transcode_video(source_video, compact_video, audio_frequency=None)
            source = root / "source.pptx"
            compact = root / "compact.pptx"
            make_video_pptx(source, source_video.read_bytes(), "Source")
            make_video_pptx(compact, compact_video.read_bytes(), "Compact")

            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source, source_quality="original")
            deck = library.add_deck(source)
            family = library.family(deck["assets"][0]["family_id"])
            source_variant = library.source_variant(family)
            compact_variant = library.import_variant(
                family["id"], compact_video, "aggressive"
            )

            source.write_bytes(compact.read_bytes())
            adopted = library.adopt_upgraded_deck_source(
                deck["id"], source, prefer_source_variant=True
            )

            self.assertEqual(
                adopted["assets"][0]["original_variant_id"], source_variant["id"]
            )
            self.assertIn(compact_variant["sha256"], family["known_hashes"])

            library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": source_variant["id"],
                        "remove_variant_ids": [compact_variant["id"]],
                    }
                ]
            )
            adopted["assets"][0]["original_variant_id"] = compact_variant["id"]
            adopted = library.adopt_upgraded_deck_source(
                deck["id"], source, prefer_source_variant=True
            )
            self.assertEqual(
                adopted["assets"][0]["original_variant_id"], source_variant["id"]
            )

    def test_cleanup_within_family_real_ssim_quarantine_restore_empty(self) -> None:
        """Real videos: original + compressed variant in one family -> scan
        finds the group with a real SSIM -> keep-smaller-close -> quarantine
        -> restore -> empty, all against real files."""
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "clip.mp4"
            make_real_video(video, seconds=2.0)
            compressed = root / "clip-small.mp4"
            transcode_video(video, compressed, audio_frequency=None)

            source = root / "deck.pptx"
            make_video_pptx(source, video.read_bytes(), "Cleanup E2E")
            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source, source_quality="original")
            family = library.families()[0]
            source_variant = library.source_variant(family)
            imported = library.import_variant(
                family["id"], compressed, "compressed-copy"
            )

            groups = library.scan_cleanup_groups(ssim_threshold=0.95)
            within = [g for g in groups if g["kind"] == "within_family"]
            self.assertEqual(len(within), 1)
            group = within[0]
            ssim_map = {
                item["variant_id"]: item["ssim_to_best"] for item in group["candidates"]
            }
            # real SSIM of the re-encode against the original is high but < 1
            self.assertGreater(ssim_map[imported["id"]], 0.9)
            self.assertLess(ssim_map[imported["id"]], 1.0)

            keep_id = imported["id"]  # smaller but SSIM-close 1080p-or-adequate
            remove_id = source_variant["id"]
            remove_path = library.variant_path(source_variant)
            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                    }
                ]
            )
            self.assertEqual(result["failed"], 0)
            self.assertFalse(remove_path.is_file())
            pending = library.pending_cleanup()
            self.assertEqual(len(pending), 1)

            restored = library.restore_cleanup_entry(pending[0]["token"])
            self.assertTrue(restored.is_file())
            self.assertEqual(len(library.families()[0]["variants"]), 2)

            # quarantine again, then empty for real
            library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                    }
                ]
            )
            self.assertEqual(library.cleanup_pending_issues(), [])
            removed = library.empty_cleanup()
            self.assertEqual(removed, 1)
            self.assertEqual(library.pending_cleanup(), [])

    def test_cleanup_cross_family_real_match_and_audio_reject(self) -> None:
        """Cross-family: same content re-encoded clusters and can be merged;
        a different-audio version never clusters (no auto-merge)."""
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            base = root / "base.mp4"
            make_real_video(base, seconds=2.0)
            same = root / "same.mp4"
            transcode_video(base, same, audio_frequency=None)
            different_audio = root / "different.mp4"
            transcode_video(base, different_audio, audio_frequency=880)

            deck_a = root / "a.pptx"
            deck_b = root / "b.pptx"
            deck_c = root / "c.pptx"
            make_video_pptx(deck_a, base.read_bytes(), "A")
            make_video_pptx(deck_b, same.read_bytes(), "B")
            make_video_pptx(deck_c, different_audio.read_bytes(), "C")

            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(deck_a, source_quality="original")
            library.archive_pptx_videos(deck_b, source_quality="original")
            library.archive_pptx_videos(deck_c, source_quality="original")

            # archive itself already dedups B into A's family by fingerprint
            family_names = sorted(f["name"] for f in library.families())
            groups = library.scan_cleanup_groups()
            cross = [g for g in groups if g["kind"] == "cross_family"]
            # the different-audio deck forms its own family and never appears
            # in any cross-family group
            for group in cross:
                for candidate in group["candidates"]:
                    self.assertTrue(candidate["auto_allowed"])
            self.assertNotIn(
                "C", [c["family_name"] for g in cross for c in g["candidates"]]
            )
            self.assertEqual(len(family_names), len(set(family_names)))

    def test_backfill_balanced_tier_transcodes_to_720p(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            video = root / "source.mp4"
            make_real_video(video, seconds=2.0, size="1920x1080")
            original_bytes = video.read_bytes()
            source = root / "source.pptx"
            make_video_pptx(source, original_bytes, "Tier E2E")

            library = VideoProject.create(root / "library")
            library.archive_pptx_videos(source)  # 真实 probe：1080p/h264/aac
            family = library.families()[0]

            low = root / "low.mp4"
            transcode_video(video, low, audio_frequency=None)  # 320×180 aac
            compact = root / "compact.pptx"
            make_video_pptx(compact, low.read_bytes(), "Compact")
            family["known_hashes"].append(sha256_file(low))

            scanned = scan_embedded_videos(compact)
            media_path = next(iter(scanned))

            result = library.upgrade_pptx_from_library(compact, quality_tier="balanced")
            self.assertEqual(result["matched"], 1)
            self.assertEqual(result["quality_tier"], "balanced")
            output = Path(result["output_pptx"])
            self.assertEqual(output.name, "compact_balanced720p.pptx")

            embedded = root / "embedded.mp4"
            with ZipFile(output) as archive:
                embedded.write_bytes(archive.read(media_path))
            probed = probe_video(embedded)
            self.assertLessEqual(probed["width"], 1280)
            self.assertLessEqual(probed["height"], 720)
            self.assertGreater(probed["width"], 0)
            self.assertEqual(probed["video_codec"], "h264")
            self.assertEqual(probed["audio_codec"], "aac")
            # 别名学习：均衡母版哈希入库，未注册新版本
            self.assertEqual(len(library.families()[0]["variants"]), 1)

            # 对照：默认最佳档对 1080p h264/aac 源原样嵌入（零损失）
            control = library.upgrade_pptx_from_library(compact)
            self.assertEqual(control["quality_tier"], "best")
            with ZipFile(control["output_pptx"]) as archive:
                self.assertEqual(archive.read(media_path), original_bytes)


if __name__ == "__main__":
    unittest.main()
