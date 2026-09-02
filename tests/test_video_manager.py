from __future__ import annotations

import copy
import hashlib
import json
import shutil
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from pptx_output_watermark.pptx_video_support import scan_embedded_videos
from pptx_tools.video_manager import (
    FAMILY_MOVE_JOURNAL_NAME,
    VideoProject,
    _is_absolute_stored_path,
    build_parser,
    _transcode_high_quality_mp4,
    normalize_library_category,
    sha256_file,
)
from scripts.batch_compact_library_pptx import (
    POLICY,
    apply_quality_selection,
    policy_for_threshold,
    compatibility_errors,
    should_replace_video,
    source_matches_report,
    target_needs_build,
)


def make_video_pptx(path: Path, video_bytes: bytes, title: str) -> None:
    video = path.with_name(f"{path.stem}.mp4")
    poster = path.with_name(f"{path.stem}.png")
    video.write_bytes(video_bytes)
    Image.new("RGB", (320, 180), "#dc6b2f").save(poster)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(5), Inches(0.5)
    ).text = title
    slide.shapes.add_movie(
        str(video),
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3.375),
        poster_frame_image=str(poster),
        mime_type="video/mp4",
    )
    presentation.save(path)


def no_probe(_: Path) -> dict[str, object]:
    return {
        "width": 320,
        "height": 180,
        "duration_sec": 2.0,
        "bitrate_kbps": 100,
        "video_codec": "h264",
        "audio_codec": "",
        "has_audio": False,
        "probe_error": "",
    }


class BackfillTranscodeTierTest(unittest.TestCase):
    def test_tier_table_shape_and_best_matches_current(self) -> None:
        from pptx_tools.video_manager import (
            BACKFILL_QUALITY_TIERS,
            DEFAULT_BACKFILL_TIER,
        )

        self.assertEqual(DEFAULT_BACKFILL_TIER, "best")
        self.assertEqual(set(BACKFILL_QUALITY_TIERS), {"best", "high", "balanced"})
        best = BACKFILL_QUALITY_TIERS["best"]
        self.assertEqual(
            (best["max_width"], best["max_height"], best["crf"]),
            (1920, 1080, 18),
        )
        self.assertEqual(best["bitrate_kbps"], 0)  # 0 = 不设限
        self.assertEqual(best["audio_bitrate"], "256k")
        self.assertEqual(best["suffix"], "high_quality")
        high = BACKFILL_QUALITY_TIERS["high"]
        self.assertEqual((high["max_width"], high["max_height"]), (1920, 1080))
        self.assertEqual((high["crf"], high["bitrate_kbps"]), (20, 12000))
        self.assertEqual(high["suffix"], "hq1080p")
        balanced = BACKFILL_QUALITY_TIERS["balanced"]
        self.assertEqual((balanced["max_width"], balanced["max_height"]), (1280, 720))
        self.assertEqual((balanced["crf"], balanced["bitrate_kbps"]), (23, 5000))
        self.assertEqual(balanced["audio_bitrate"], "128k")
        self.assertEqual(balanced["suffix"], "balanced720p")

    def _capture_command(self, audio_usable: bool = True) -> list[str]:
        from pptx_tools import video_manager as vm

        with (
            patch.object(vm, "_audio_stream_usable", return_value=audio_usable),
            patch.object(vm, "run_binary") as run,
        ):
            vm._transcode_high_quality_mp4(Path("in.mp4"), Path("out.mp4"), 3840, 2160)
        return list(run.call_args.args[0])

    def test_default_command_matches_current_plus_fps_passthrough(self) -> None:
        self.assertEqual(
            self._capture_command(),
            [
                "ffmpeg",
                "-y",
                "-i",
                "in.mp4",
                "-map",
                "0:v:0",
                "-map",
                "0:a?",
                "-vf",
                "scale=1920:1080:flags=lanczos",
                "-fps_mode",
                "passthrough",
                "-c:v",
                "libx264",
                "-preset",
                "medium",
                "-crf",
                "18",
                "-profile:v",
                "main",
                "-bf",
                "0",
                "-pix_fmt",
                "yuv420p",
                "-c:a",
                "aac",
                "-b:a",
                "256k",
                "-movflags",
                "+faststart",
                "out.mp4",
            ],
        )

    def test_tier_command_applies_crf_maxrate_audio_and_scale(self) -> None:
        from pptx_tools import video_manager as vm

        with (
            patch.object(vm, "_audio_stream_usable", return_value=True),
            patch.object(vm, "run_binary") as run,
        ):
            vm._transcode_high_quality_mp4(
                Path("in.mp4"),
                Path("out.mp4"),
                1920,
                1080,
                max_width=1280,
                max_height=720,
                crf=23,
                bitrate_kbps=5000,
                audio_bitrate="128k",
            )
        cmd = list(run.call_args.args[0])
        self.assertIn("scale=1280:720:flags=lanczos", cmd)
        self.assertEqual(cmd[cmd.index("-crf") + 1], "23")
        self.assertEqual(cmd[cmd.index("-maxrate") + 1], "5000k")
        self.assertEqual(cmd[cmd.index("-bufsize") + 1], "10000k")
        self.assertEqual(cmd[cmd.index("-b:a") + 1], "128k")
        # 默认档不出现码率上限
        self.assertNotIn("-maxrate", self._capture_command())

    def test_empty_audio_track_drops_audio_arguments(self) -> None:
        cmd = self._capture_command(audio_usable=False)
        self.assertNotIn("0:a?", cmd)
        self.assertNotIn("-c:a", cmd)
        self.assertNotIn("-b:a", cmd)

    def test_audio_stream_usable_only_false_when_proven_empty(self) -> None:
        from pptx_tools import video_manager as vm

        class Result:
            def __init__(self, payload: dict) -> None:
                self.stdout = json.dumps(payload)

        cases = [
            ({"streams": []}, True),  # 无音轨：保持可选映射（现状等价）
            ({"streams": [{}]}, True),  # 字段未知
            ({"streams": [{"duration": "0.000000", "nb_frames": "1"}]}, False),
            ({"streams": [{"duration": "0.000000", "nb_frames": "120"}]}, True),
            ({"streams": [{"duration": "2.5", "nb_frames": "0"}]}, True),
            ({"streams": [{"duration": "N/A", "nb_frames": "0"}]}, True),
        ]
        for payload, expected in cases:
            with patch.object(
                vm, "run_binary", side_effect=lambda *a, _p=payload, **k: Result(_p)
            ):
                self.assertIs(vm._audio_stream_usable(Path("x.mp4")), expected, payload)
        with patch.object(vm, "run_binary", side_effect=RuntimeError("boom")):
            self.assertIs(vm._audio_stream_usable(Path("x.mp4")), True)


class StoredPathPortabilityTest(unittest.TestCase):
    def test_absolute_paths_are_recognized_across_platforms(self) -> None:
        self.assertTrue(_is_absolute_stored_path("/Volumes/example/source.pptx"))
        self.assertTrue(_is_absolute_stored_path(r"X:\Portable\source.pptx"))
        self.assertTrue(_is_absolute_stored_path(r"\\server\share\source.pptx"))
        self.assertTrue(_is_absolute_stored_path("~/source.pptx"))
        self.assertFalse(_is_absolute_stored_path("../outside/source.pptx"))
        self.assertFalse(_is_absolute_stored_path("media/source.mp4"))


class BackfillCompatibilityTest(unittest.TestCase):
    def _meta(self, **overrides):
        metadata = {
            "suffix": ".mp4",
            "video_codec": "h264",
            "audio_codec": "aac",
            "width": 1920,
            "height": 1080,
            "bitrate_kbps": 8000,
        }
        metadata.update(overrides)
        return metadata

    def test_best_tier_matches_current_semantics(self) -> None:
        from pptx_tools.video_manager import _backfill_compatibility, _tier_spec

        spec = _tier_spec("best")
        self.assertTrue(_backfill_compatibility(self._meta(), spec))
        self.assertTrue(_backfill_compatibility(self._meta(bitrate_kbps=50000), spec))
        self.assertFalse(_backfill_compatibility(self._meta(suffix=".wmv"), spec))
        self.assertFalse(_backfill_compatibility(self._meta(video_codec="hevc"), spec))
        self.assertFalse(_backfill_compatibility(self._meta(audio_codec="ac3"), spec))
        self.assertFalse(
            _backfill_compatibility(self._meta(width=3840, height=2160), spec)
        )
        self.assertFalse(_backfill_compatibility(self._meta(width=1919), spec))  # 奇数
        self.assertFalse(_backfill_compatibility(self._meta(width=0), spec))

    def test_capped_tiers_add_resolution_and_bitrate_ceilings(self) -> None:
        from pptx_tools.video_manager import _backfill_compatibility, _tier_spec

        balanced = _tier_spec("balanced")
        high = _tier_spec("high")
        self.assertFalse(
            _backfill_compatibility(self._meta(), balanced)
        )  # 1080p 超 720p
        self.assertTrue(
            _backfill_compatibility(
                self._meta(width=1280, height=720, bitrate_kbps=4000), balanced
            )
        )
        self.assertFalse(
            _backfill_compatibility(
                self._meta(width=1280, height=720, bitrate_kbps=8000), balanced
            )
        )
        self.assertTrue(_backfill_compatibility(self._meta(bitrate_kbps=12000), high))
        self.assertFalse(_backfill_compatibility(self._meta(bitrate_kbps=20000), high))
        # 码率探测缺失：非 best 档视为超限
        self.assertFalse(
            _backfill_compatibility(
                self._meta(width=1280, height=720, bitrate_kbps=0), balanced
            )
        )
        # 竖版包络：均衡档 720×1280
        self.assertTrue(
            _backfill_compatibility(
                self._meta(width=720, height=1280, bitrate_kbps=4000), balanced
            )
        )
        self.assertFalse(
            _backfill_compatibility(
                self._meta(width=1080, height=1920, bitrate_kbps=4000), balanced
            )
        )

    def test_plan_backfill_action_texts(self) -> None:
        from pptx_tools.video_manager import plan_backfill_action

        self.assertEqual(
            plan_backfill_action(
                self._meta(width=1280, height=720, bitrate_kbps=4000), "balanced"
            ),
            "原样嵌入（已达标）",
        )
        self.assertEqual(
            plan_backfill_action(self._meta(), "balanced"),
            "转码至 ≤720p · CRF 23 · ≤5Mbps",
        )
        self.assertEqual(
            plan_backfill_action(self._meta(), "best"), "原样嵌入（已达标）"
        )
        self.assertEqual(
            plan_backfill_action(self._meta(width=0, height=0), "balanced"),
            "按档位规格回填",
        )
        with self.assertRaises(ValueError):
            plan_backfill_action(self._meta(), "bogus")

    def test_plan_matches_delivery_decision(self) -> None:
        # 同一 metadata：预览分支与共享判定一致
        from pptx_tools.video_manager import (
            _backfill_compatibility,
            _tier_spec,
            plan_backfill_action,
        )

        for meta, tier in [
            (self._meta(), "best"),
            (self._meta(), "balanced"),
            (self._meta(width=1280, height=720, bitrate_kbps=4000), "balanced"),
            (self._meta(bitrate_kbps=20000), "high"),
        ]:
            compatible = _backfill_compatibility(meta, _tier_spec(tier))
            text = plan_backfill_action(meta, tier)
            self.assertEqual(compatible, text == "原样嵌入（已达标）", (meta, tier))


class DeliveryMasterTierTest(unittest.TestCase):
    def _library_with_source(self, root: Path, metadata: dict) -> tuple:
        source = root / "source.pptx"
        make_video_pptx(source, b"tier-source", "Source")
        library = VideoProject.create(root / "library")
        with patch("pptx_tools.video_manager.probe_video", return_value=metadata):
            library.archive_pptx_videos(source)
        family = library.families()[0]
        return library, family

    def _base_metadata(self) -> dict:
        return {
            **no_probe(Path()),
            "width": 1920,
            "height": 1080,
            "video_codec": "h264",
            "audio_codec": "aac",
            "has_audio": True,
            "bitrate_kbps": 8000,
        }

    def test_default_tier_embeds_compatible_source_as_is(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            library, family = self._library_with_source(root, self._base_metadata())
            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4"
            ) as transcode:
                delivery, _ = library._delivery_master(family, root)
            transcode.assert_not_called()
            self.assertEqual(delivery.suffix.lower(), ".mp4")

    def test_balanced_tier_transcodes_1080p_source_with_tier_params(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            library, family = self._library_with_source(root, self._base_metadata())

            def fake_transcode(_source, target, _w, _h, **kwargs):
                self.assertEqual(kwargs["crf"], 23)
                self.assertEqual(kwargs["bitrate_kbps"], 5000)
                self.assertEqual(kwargs["audio_bitrate"], "128k")
                self.assertEqual(
                    (kwargs["max_width"], kwargs["max_height"]), (1280, 720)
                )
                target.write_bytes(b"balanced-delivery")

            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4",
                side_effect=fake_transcode,
            ):
                delivery, digest = library._delivery_master(
                    family, root, tier="balanced"
                )
            self.assertEqual(delivery, root / f"{family['id']}.mp4")
            self.assertEqual(digest, hashlib.sha256(b"balanced-delivery").hexdigest())

    def test_balanced_tier_embeds_small_low_bitrate_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            metadata = self._base_metadata()
            metadata.update({"width": 1280, "height": 720, "bitrate_kbps": 4000})
            library, family = self._library_with_source(root, metadata)
            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4"
            ) as transcode:
                delivery, digest = library._delivery_master(
                    family, root, tier="balanced"
                )
            transcode.assert_not_called()
            self.assertEqual(digest, library.source_variant(family)["sha256"])

    def test_high_tier_transcodes_high_bitrate_1080p(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            metadata = self._base_metadata()
            metadata["bitrate_kbps"] = 20000
            library, family = self._library_with_source(root, metadata)

            def fake_transcode(_source, target, _w, _h, **kwargs):
                self.assertEqual(kwargs["bitrate_kbps"], 12000)
                target.write_bytes(b"hq-delivery")

            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4",
                side_effect=fake_transcode,
            ):
                library._delivery_master(family, root, tier="high")

    def test_invalid_tier_raises(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            library, family = self._library_with_source(root, self._base_metadata())
            with self.assertRaises(ValueError):
                library._delivery_master(family, root, tier="bogus")


class UpgradeQualityTierTest(unittest.TestCase):
    def _setup_library(self, root: Path) -> tuple:
        source = root / "source.pptx"
        compact = root / "compact.pptx"
        make_video_pptx(source, b"tier-source-bytes", "Source")
        make_video_pptx(compact, b"tier-compact-bytes", "Compact")
        library = VideoProject.create(root / "library")
        metadata = {
            **no_probe(Path()),
            "width": 1920,
            "height": 1080,
            "video_codec": "h264",
            "audio_codec": "aac",
            "has_audio": True,
            "bitrate_kbps": 8000,
        }
        with patch("pptx_tools.video_manager.probe_video", return_value=metadata):
            library.archive_pptx_videos(source)
        family = library.families()[0]
        family["known_hashes"].append(sha256_file(compact.with_suffix(".mp4")))
        return library, compact, family

    def test_default_output_name_unchanged(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library, compact, _ = self._setup_library(Path(temp_dir))
            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4",
                side_effect=lambda _s, target, _w, _h, **_k: target.write_bytes(b"x"),
            ):
                result = library.upgrade_pptx_from_library(compact)
            self.assertEqual(result["output_pptx"].name, "compact_high_quality.pptx")
            self.assertEqual(result["quality_tier"], "best")

    def test_tier_suffix_and_result_field(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library, compact, _ = self._setup_library(Path(temp_dir))

            def fake_transcode(_s, target, _w, _h, **kwargs):
                self.assertEqual(kwargs["crf"], 23)
                target.write_bytes(b"balanced-master")

            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4",
                side_effect=fake_transcode,
            ):
                result = library.upgrade_pptx_from_library(
                    compact, quality_tier="balanced"
                )
            self.assertEqual(result["output_pptx"].name, "compact_balanced720p.pptx")
            self.assertEqual(result["quality_tier"], "balanced")
            # 别名学习仍生效：转码母版哈希入 known_hashes
            digest = hashlib.sha256(b"balanced-master").hexdigest()
            self.assertIn(digest, library.families()[0]["known_hashes"])

    def test_invalid_tier_raises_before_writing(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library, compact, _ = self._setup_library(Path(temp_dir))
            with self.assertRaises(ValueError):
                library.upgrade_pptx_from_library(compact, quality_tier="bogus")
            self.assertFalse((Path(temp_dir) / "compact_bogus.pptx").exists())


class VideoProjectTest(unittest.TestCase):
    def test_upgrade_deduplicates_identical_delivery_media_parts(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            compact = root / "compact.pptx"
            make_video_pptx(source, b"master-video", "Master")

            poster = root / "poster.png"
            first = root / "first.mp4"
            second = root / "second.mp4"
            Image.new("RGB", (320, 180), "#dc6b2f").save(poster)
            first.write_bytes(b"compact-one")
            second.write_bytes(b"compact-two")
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_movie(
                str(first), Inches(1), Inches(1), Inches(2), Inches(1), str(poster)
            )
            slide.shapes.add_movie(
                str(second), Inches(4), Inches(1), Inches(2), Inches(1), str(poster)
            )
            presentation.save(compact)

            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(source)
            family = library.families()[0]
            with ZipFile(compact) as archive:
                compact_hashes = [
                    hashlib.sha256(archive.read(name)).hexdigest()
                    for name in sorted(
                        item
                        for item in archive.namelist()
                        if item.startswith("ppt/media/") and item.endswith(".mp4")
                    )
                ]
            family["known_hashes"].extend(compact_hashes)

            result = library.upgrade_pptx_from_library(compact)

            output = Path(result["output_pptx"])
            scanned = scan_embedded_videos(output)
            self.assertEqual(len(scanned), 1)
            self.assertEqual(sum(len(item.occurrences) for item in scanned.values()), 2)
            with ZipFile(output) as archive:
                self.assertEqual(
                    len(
                        [
                            name
                            for name in archive.namelist()
                            if name.startswith("ppt/media/") and name.endswith(".mp4")
                        ]
                    ),
                    1,
                )

    def test_batch_reuse_requires_unchanged_source(self) -> None:
        target = {"source_sha256": "same", "source_size_bytes": 100}
        self.assertTrue(source_matches_report(target, dict(target)))
        self.assertFalse(
            source_matches_report(target, {**target, "source_sha256": "changed"})
        )
        self.assertFalse(
            source_matches_report(target, {**target, "source_size_bytes": 101})
        )
        self.assertFalse(source_matches_report(target, None))

    def test_batch_only_rebuilds_changed_or_missing_outputs(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            output = Path(temporary) / "output.pptx"
            output.write_bytes(b"ok")
            target = {"source_sha256": "source", "source_size_bytes": 100}
            existing = {
                **target,
                "policy": POLICY,
                "status": "validated",
                "output_sha256": sha256_file(output),
            }
            self.assertFalse(target_needs_build(target, existing, output))
            self.assertTrue(
                target_needs_build(
                    {**target, "source_sha256": "changed"}, existing, output
                )
            )
            output.unlink()
            self.assertTrue(target_needs_build(target, existing, output))
            self.assertFalse(
                target_needs_build(target, {**existing, "status": "skipped"}, output)
            )

    def test_batch_quality_selection_requires_matching_profile_and_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Quality")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            original = project.source_variant(family)

            def fake_compress(args, **_: object) -> dict[str, object]:
                Path(args.output).write_bytes(b"quality-video")
                return {"output_pptx": Path(args.output)}

            with (
                patch(
                    "pptx_tools.video_manager.compact_standalone_video",
                    side_effect=fake_compress,
                ),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                selected = project.compress_variant(
                    original["id"], "balanced", activate=False
                )

            decision = {
                "source_sha256": original["sha256"],
                "selected_profile": "balanced",
                "selected_variant_id": selected["id"],
                "selected_sha256": selected["sha256"],
                "ssim": 0.97,
            }
            report = {"families": {family["id"]: {}}}
            apply_quality_selection(
                project,
                {family["id"]: selected},
                {},
                {"threshold": 0.95, "items": {family["id"]: decision}},
                root / "delivery",
                report,
            )
            self.assertEqual(
                report["families"][family["id"]]["selected_variant_id"],
                selected["id"],
            )

            source_report = {"families": {family["id"]: {}}}
            source_blocked: dict[str, str] = {}
            apply_quality_selection(
                project,
                {family["id"]: selected},
                source_blocked,
                {
                    "threshold": 0.95,
                    "items": {
                        family["id"]: {
                            **decision,
                            "selected_variant_id": original["id"],
                            "selected_sha256": original["sha256"],
                            "ssim": 1.0,
                        }
                    },
                },
                root / "delivery",
                source_report,
            )
            self.assertEqual(
                source_report["families"][family["id"]]["selected_profile"],
                "source",
            )
            self.assertIn(family["id"], source_blocked)

            self.assertEqual(
                policy_for_threshold(0.90),
                "ssim090-video-15pct-and-pptx-5pct-or-10mib-v4",
            )
            apply_quality_selection(
                project,
                {family["id"]: selected},
                {},
                {
                    "threshold": 0.90,
                    "items": {family["id"]: {**decision, "ssim": 0.90}},
                },
                root / "delivery-090",
                {"families": {family["id"]: {}}},
            )
            with self.assertRaises(RuntimeError):
                apply_quality_selection(
                    project,
                    {family["id"]: selected},
                    {},
                    {
                        "threshold": 0.89,
                        "items": {family["id"]: decision},
                    },
                    root / "delivery-089",
                    {"families": {family["id"]: {}}},
                )

            for field, value in (
                ("selected_profile", "high"),
                ("source_sha256", "wrong-source"),
            ):
                invalid = {**decision, field: value}
                with self.subTest(field=field), self.assertRaises(RuntimeError):
                    apply_quality_selection(
                        project,
                        {family["id"]: selected},
                        {},
                        {"threshold": 0.95, "items": {family["id"]: invalid}},
                        root / "delivery",
                        {"families": {family["id"]: {}}},
                    )

    def test_batch_delivery_blocks_timing_and_audio_changes(self) -> None:
        source = {
            "duration_sec": 60.0,
            "width": 1920,
            "height": 1080,
            "has_audio": True,
        }
        compatible = {
            **source,
            "duration_sec": 60.1,
            "width": 1280,
            "height": 720,
            "video_codec": "h264",
            "audio_codec": "aac",
        }
        changed = {**compatible, "duration_sec": 60.5, "has_audio": False}

        self.assertEqual(compatibility_errors(source, compatible), [])
        self.assertEqual(len(compatibility_errors(source, changed)), 2)

    def test_batch_delivery_requires_meaningful_lossy_savings(self) -> None:
        self.assertFalse(
            should_replace_video("ppt/media/video.mp4", 1000, 900, compatible=True)
        )
        self.assertTrue(
            should_replace_video("ppt/media/video.mp4", 1000, 800, compatible=True)
        )
        self.assertTrue(
            should_replace_video("ppt/media/video.wmv", 1000, 990, compatible=True)
        )
        self.assertFalse(
            should_replace_video("ppt/media/video.wmv", 1000, 500, compatible=False)
        )

    def test_agent_cli_exposes_video_import_and_source_selection(self) -> None:
        imported = build_parser().parse_args(
            [
                "import-video",
                "/tmp/library",
                "/tmp/video.mp4",
                "--family-id",
                "family-1",
                "--source-quality",
                "original",
            ]
        )
        selected = build_parser().parse_args(
            ["set-source", "/tmp/library", "variant-1"]
        )

        self.assertEqual(imported.action, "import-video")
        self.assertEqual(imported.family_id, "family-1")
        self.assertEqual(imported.source_quality, "original")
        self.assertEqual(selected.action, "set-source")
        self.assertEqual(selected.variant_id, "variant-1")

    def test_high_quality_mp4_disables_b_frames_for_sparse_wmv_timestamps(self) -> None:
        with patch("pptx_tools.video_manager.run_binary") as run:
            _transcode_high_quality_mp4(
                Path("source.wmv"), Path("target.mp4"), 1680, 1048
            )

        command = run.call_args.args[0]
        self.assertEqual(command[command.index("-bf") + 1], "0")

    def test_library_category_is_relative_and_nested(self) -> None:
        self.assertEqual(
            normalize_library_category("示例项目/2026/外发/"),
            Path("示例项目/2026/外发"),
        )
        for value in ("../outside", "/absolute", "C:\\outside", "a//b"):
            with self.subTest(value=value), self.assertRaises(ValueError):
                normalize_library_category(value)

    def test_external_import_uses_shared_name_and_category_rules(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "1662562042-示例设备_fault.mp4"
            source.write_bytes(b"video")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 1000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.import_external_video(
                    source, source_quality="original", category="示例分类/异常样本"
                )

            family = library.families()[0]
            self.assertEqual(family["name"], "示例设备_异常_1662562042")
            self.assertEqual(family["category"], "示例分类/异常样本")

    def test_exact_duplicates_are_archived_once_across_decks(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"same-video", "First")
            make_video_pptx(second, b"same-video", "Second")
            project = VideoProject.create(root / "project")

            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(first)
                project.add_deck(second)

            self.assertEqual(len(project.decks()), 2)
            self.assertEqual(len(project.families()), 1)
            self.assertEqual(len(project.families()[0]["variants"]), 1)
            self.assertEqual(
                project.decks()[0]["assets"][0]["family_id"],
                project.decks()[1]["assets"][0]["family_id"],
            )

    def test_identical_deck_copy_is_recorded_as_source_alias(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            alias = root / "renamed" / "copy.pptx"
            make_video_pptx(source, b"same-deck-video", "Source")
            alias.parent.mkdir()
            shutil.copy2(source, alias)
            project = VideoProject.create(root / "project")

            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
                project.add_deck(alias)

            self.assertEqual(len(project.decks()), 1)
            self.assertEqual(
                project.decks()[0]["source_aliases"],
                [str(alias.resolve())],
            )
            self.assertEqual(len(project.families()), 1)

    def test_library_archive_deduplicates_without_tracking_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"shared-source-video", "First")
            make_video_pptx(second, b"shared-source-video", "Second")
            library = VideoProject.create(root / "library")

            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                first_result = library.archive_pptx_videos(first, category="示例项目")
                second_result = library.archive_pptx_videos(second, category="客户B")

            self.assertEqual(first_result["added"], 1)
            self.assertEqual(second_result["reused"], 1)
            self.assertEqual(len(library.families()), 1)
            self.assertEqual(library.decks(), [])
            self.assertEqual(len(list((library.root / "media").rglob("*.mp4"))), 1)
            self.assertFalse((library.root / "media" / "客户B").exists())

    def test_deck_reuses_archived_compressed_hash_alias(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            compressed = root / "compressed.pptx"
            make_video_pptx(source, b"high-quality-source", "Source")
            make_video_pptx(compressed, b"compressed-alias", "Compressed")
            project = VideoProject.create(root / "project")

            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.archive_pptx_videos(source)
                family = project.families()[0]
                family.setdefault("known_hashes", []).append(
                    hashlib.sha256(b"compressed-alias").hexdigest()
                )
                project.save()
                deck = project.add_deck(compressed)

            self.assertEqual(len(project.families()), 1)
            self.assertEqual(deck["assets"][0]["family_id"], family["id"])

    def test_reencoded_video_in_deck_matches_by_content_fingerprint(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            reencoded = root / "reencoded.pptx"
            make_video_pptx(source, b"source-encoding", "Source")
            make_video_pptx(reencoded, b"different-encoding", "Reencoded")
            project = VideoProject.create(root / "project")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [128] * 5,
                "has_audio": False,
            }

            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                first = project.add_deck(source)
                second = project.add_deck(reencoded)

            self.assertEqual(len(project.families()), 1)
            self.assertEqual(
                first["assets"][0]["family_id"], second["assets"][0]["family_id"]
            )
            self.assertEqual(len(project.families()[0]["known_hashes"]), 2)

    def test_damaged_video_prefix_matches_unique_named_complete_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            project = VideoProject.create(root / "project")
            complete = project.root / "media" / "production-line.mp4"
            prefix = (b"video-payload-" * 90_000)[: 1024 * 1024]
            complete.write_bytes(prefix + b"complete-tail")
            digest = sha256_file(complete)
            variant = {
                "id": "complete-variant",
                "label": "source",
                "profile": "original",
                "path": project.encode_path(complete),
                "sha256": digest,
                "size_bytes": complete.stat().st_size,
                "mtime_ns": complete.stat().st_mtime_ns,
                "created_at": "2026-07-22T00:00:00+00:00",
                "source_variant_id": None,
                **no_probe(complete),
            }
            family = {
                "id": "complete-family",
                "name": "示例生产线",
                "category": "",
                "source_variant_id": variant["id"],
                "active_variant_id": variant["id"],
                "known_hashes": [digest],
                "source_hashes": [digest],
                "variants": [variant],
            }
            project.families().append(family)
            damaged = root / "damaged.mp4"
            damaged.write_bytes(prefix)
            damaged_metadata = {
                **no_probe(damaged),
                "width": 0,
                "height": 0,
                "duration_sec": 0,
                "probe_error": "moov atom not found",
            }

            self.assertIs(
                project._family_by_damaged_prefix(
                    "示例生产线",
                    damaged,
                    sha256_file(damaged),
                    damaged_metadata,
                ),
                family,
            )
            self.assertIsNone(
                project._family_by_damaged_prefix(
                    "另一个视频",
                    damaged,
                    sha256_file(damaged),
                    damaged_metadata,
                )
            )

            deck = root / "damaged.pptx"
            make_video_pptx(deck, prefix, "Damaged")
            presentation = Presentation(deck)
            presentation.slides[0].shapes[-1].name = "示例生产线"
            presentation.save(deck)
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    return_value=damaged_metadata,
                ),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                result = project.archive_pptx_videos(deck)

            self.assertEqual(result["reused"], 1)
            self.assertEqual(result["added"], 0)
            self.assertEqual(len(project.families()), 1)
            self.assertIn(hashlib.sha256(prefix).hexdigest(), family["known_hashes"])

    def test_archive_uses_category_and_family_rename_renames_source_file(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"rename-stable-video", "Original name")
            library = VideoProject.create(root / "library")
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                library.archive_pptx_videos(source, category="示例项目/2026")

            family = library.families()[0]
            family_id = family["id"]
            variant = library.source_variant(family)
            digest = variant["sha256"]
            original_path = library.variant_path(variant)
            self.assertEqual(
                original_path.parent, library.root / "media" / "示例项目" / "2026"
            )

            renamed = library.rename_family_and_source(family_id, "发布会主视频")

            self.assertFalse(original_path.exists())
            self.assertEqual(
                renamed.name, f"发布会主视频_[320x180_2.0s]_{digest[:8]}.mp4"
            )
            self.assertEqual(family["name"], "发布会主视频")
            self.assertEqual(family["id"], family_id)
            self.assertEqual(sha256_file(renamed), digest)
            self.assertIs(library.family_by_known_hash(digest), family)

            dotted = library.rename_family_and_source(family_id, "5.2追溯")
            self.assertEqual(dotted.name, f"5.2追溯_[320x180_2.0s]_{digest[:8]}.mp4")
            self.assertEqual(family["name"], "5.2追溯")

    def test_family_source_rename_rolls_back_on_manifest_conflict(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"rollback-family-rename", "Original")
            library = VideoProject.create(root / "library")
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                library.archive_pptx_videos(source)
            family = library.families()[0]
            original_name = family["name"]
            original_path = library.variant_path(library.source_variant(family))

            with patch.object(library, "save", side_effect=RuntimeError("conflict")):
                with self.assertRaisesRegex(RuntimeError, "conflict"):
                    library.rename_family_and_source(family["id"], "new-name")

            self.assertEqual(family["name"], original_name)
            self.assertTrue(original_path.is_file())
            self.assertFalse(original_path.with_name("new-name.mp4").exists())

    def test_1080p_archive_keeps_original_hash_as_alias(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            original_bytes = b"four-k-source"
            make_video_pptx(source, original_bytes, "4K")
            library = VideoProject.create(root / "library")
            high = {**no_probe(Path()), "width": 3840, "height": 2160}
            archived = {**no_probe(Path()), "width": 1920, "height": 1080}

            def fake_transcode(
                _source: Path,
                target: Path,
                _width: int,
                _height: int,
                **_kwargs: object,
            ) -> None:
                target.write_bytes(b"high-quality-1080p")

            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[high, archived],
                ),
                patch(
                    "pptx_tools.video_manager._transcode_high_quality_mp4",
                    side_effect=fake_transcode,
                ),
            ):
                library.archive_pptx_videos(source, source_quality="1080p")

            family = library.families()[0]
            variant = library.source_variant(family)
            self.assertEqual(variant["profile"], "1080p_source")
            self.assertEqual(
                library.variant_path(variant).read_bytes(), b"high-quality-1080p"
            )
            self.assertIn(
                sha256_file(source.with_suffix(".mp4")), family["known_hashes"]
            )
            self.assertIn(variant["sha256"], family["known_hashes"])
            self.assertIn(
                sha256_file(source.with_suffix(".mp4")), family["source_hashes"]
            )

    def test_mp4_import_keeps_resolution_and_original_hash_alias(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.wmv"
            source.write_bytes(b"wmv-source")
            library = VideoProject.create(root / "library")
            original = {**no_probe(Path()), "width": 2560, "height": 1440}
            normalized = {**no_probe(Path()), "width": 2560, "height": 1440}

            def fake_transcode(
                _source: Path,
                target: Path,
                _width: int,
                _height: int,
                **kwargs: object,
            ) -> None:
                self.assertFalse(kwargs["limit_1080p"])
                target.write_bytes(b"compatible-mp4")

            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[original, normalized],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value={
                        "duration_ms": 1000,
                        "aspect_ppm": 1_777_778,
                        "frames": ["0" * 16] * 5,
                        "luma": [120] * 5,
                        "has_audio": False,
                    },
                ),
                patch(
                    "pptx_tools.video_manager._transcode_high_quality_mp4",
                    side_effect=fake_transcode,
                ),
            ):
                result = library.import_external_video(source, source_quality="mp4")

            family = library.families()[0]
            variant = library.source_variant(family)
            self.assertEqual(result["status"], "created")
            self.assertEqual(variant["profile"], "mp4_source")
            self.assertEqual(variant["width"], 2560)
            self.assertEqual(variant["height"], 1440)
            self.assertEqual(library.variant_path(variant).suffix, ".mp4")
            self.assertIn(sha256_file(source), family["known_hashes"])
            self.assertIn(variant["sha256"], family["known_hashes"])
            self.assertEqual(variant["origin_paths"], [str(source.resolve())])

    def test_upgrade_skips_video_already_at_library_source_quality(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-quality-video", "Source")
            library = VideoProject.create(root / "library")
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                library.archive_pptx_videos(source)
                result = library.upgrade_pptx_from_library(source)

            self.assertIsNone(result["output_pptx"])
            self.assertEqual(result["already_high_quality"], 1)
            self.assertEqual(result["matched"], 0)

    def test_delivery_master_reuses_compatible_m4v_without_transcoding(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"compatible-m4v", "M4V")
            library = VideoProject.create(root / "library")
            metadata = {
                **no_probe(Path()),
                "width": 1920,
                "height": 1080,
                "audio_codec": "aac",
                "has_audio": True,
            }
            with patch("pptx_tools.video_manager.probe_video", return_value=metadata):
                library.archive_pptx_videos(source)
            family = library.families()[0]
            variant = library.source_variant(family)
            managed = library.variant_path(variant)
            m4v = managed.with_suffix(".m4v")
            managed.rename(m4v)
            variant["path"] = library.encode_path(m4v)

            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4"
            ) as transcode:
                delivery, digest = library._delivery_master(family, root)

            transcode.assert_not_called()
            self.assertEqual(delivery, m4v)
            self.assertEqual(digest, variant["sha256"])

    def test_upgrade_registers_transcoded_delivery_hash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            compact = root / "compact.pptx"
            make_video_pptx(source, b"large-source", "Source")
            make_video_pptx(compact, b"small-proxy", "Compact")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(source)
            family = library.families()[0]
            source_variant = library.source_variant(family)
            source_variant.update({"width": 3840, "height": 2160})
            proxy_digest = sha256_file(compact.with_suffix(".mp4"))
            family["known_hashes"].append(proxy_digest)

            def fake_transcode(
                _source: Path,
                target: Path,
                _width: int,
                _height: int,
                **_kwargs: object,
            ) -> None:
                target.write_bytes(b"compatible-delivery")

            with patch(
                "pptx_tools.video_manager._transcode_high_quality_mp4",
                side_effect=fake_transcode,
            ):
                result = library.upgrade_pptx_from_library(compact)

            delivery_digest = hashlib.sha256(b"compatible-delivery").hexdigest()
            self.assertEqual(result["aliases_added"], 1)
            self.assertIn(delivery_digest, family["known_hashes"])

    def test_later_higher_resolution_source_requires_manual_promotion(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"same-content-720p", "First")
            make_video_pptx(second, b"same-content-1080p", "Second")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            low = {**no_probe(Path()), "width": 1280, "height": 720}
            high = {**no_probe(Path()), "width": 1920, "height": 1080}
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[low, high],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.archive_pptx_videos(first, source_quality="original")
                library.archive_pptx_videos(second, source_quality="original")

            family = library.families()[0]
            self.assertEqual(len(library.families()), 1)
            self.assertEqual(len(family["variants"]), 2)
            master = library.source_variant(family)
            self.assertEqual((master["width"], master["height"]), (1280, 720))
            self.assertEqual(
                library.variant_path(master).read_bytes(), b"same-content-720p"
            )

            high_variant = next(
                item for item in family["variants"] if item["id"] != master["id"]
            )
            library.set_source_variant(high_variant["id"])
            self.assertEqual(
                library.variant_path(library.source_variant(family)).read_bytes(),
                b"same-content-1080p",
            )

            upgraded = library.upgrade_pptx_from_library(first)
            with ZipFile(upgraded["output_pptx"]) as archive:
                media_path = next(iter(scan_embedded_videos(first)))
                self.assertEqual(archive.read(media_path), b"same-content-1080p")

    def test_same_resolution_higher_bitrate_is_retained_as_candidate(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"low-bitrate", "First")
            make_video_pptx(second, b"higher-bitrate", "Second")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            low = {
                **no_probe(Path()),
                "width": 1920,
                "height": 1080,
                "bitrate_kbps": 1000,
            }
            high = {**low, "bitrate_kbps": 1600}
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[low, high],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.archive_pptx_videos(first, source_quality="original")
                result = library.archive_pptx_videos(second, source_quality="original")

            family = library.families()[0]
            self.assertEqual(result["candidates_added"], 1)
            self.assertEqual(len(family["variants"]), 2)
            self.assertEqual(
                library.variant_path(library.source_variant(family)).read_bytes(),
                b"low-bitrate",
            )

    def test_lower_resolution_reencode_never_replaces_family_master(self) -> None:
        """A smaller compressed re-encode must not downgrade the library
        source: the upgrade comparison uses the candidate's effective stored
        size (after the optional 1080p downscale), not the envelope cap."""
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"master-640x360", "First")
            make_video_pptx(second, b"compressed-320x180", "Second")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            high = {**no_probe(Path()), "width": 640, "height": 360}
            low = {**no_probe(Path()), "width": 320, "height": 180}
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[high, low],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.archive_pptx_videos(first, source_quality="1080p")
                library.archive_pptx_videos(second, source_quality="1080p")

            self.assertEqual(len(library.families()), 1)
            family = library.families()[0]
            # no new variant, no churn: the low-quality re-encode only
            # registers its hash alias
            self.assertEqual(len(family["variants"]), 1)
            master = library.source_variant(family)
            self.assertEqual((master["width"], master["height"]), (640, 360))
            self.assertEqual(
                library.variant_path(master).read_bytes(), b"master-640x360"
            )

    def test_same_resolution_reencode_does_not_churn_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"master-bytes", "First")
            make_video_pptx(second, b"reencoded-bytes", "Second")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            same = {**no_probe(Path()), "width": 640, "height": 360}
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[dict(same), dict(same)],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.archive_pptx_videos(first, source_quality="1080p")
                library.archive_pptx_videos(second, source_quality="1080p")

            family = library.families()[0]
            self.assertEqual(len(family["variants"]), 1)
            self.assertEqual(
                library.variant_path(library.source_variant(family)).read_bytes(),
                b"master-bytes",
            )

    def test_import_rejects_exact_copy_already_in_another_family(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            duplicate = root / "renamed-copy.mp4"
            make_video_pptx(first, b"first-source", "First")
            make_video_pptx(second, b"second-source", "Second")
            duplicate.write_bytes(b"first-source")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(first)
                library.archive_pptx_videos(second)
                second_family = library.families()[1]
                with self.assertRaisesRegex(ValueError, "already exists"):
                    library.import_variant(second_family["id"], duplicate)

            self.assertEqual(len(library.families()), 2)
            self.assertEqual(len(list((library.root / "media").rglob("*.mp4"))), 2)

    def test_external_video_matches_family_without_automatic_promotion(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck = root / "deck.pptx"
            external = root / "external.mp4"
            make_video_pptx(deck, b"low-source", "Demo")
            external.write_bytes(b"high-source")
            library = VideoProject.create(root / "library")
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            low = {**no_probe(Path()), "width": 1280, "height": 720}
            high = {**no_probe(Path()), "width": 1920, "height": 1080}
            with (
                patch(
                    "pptx_tools.video_manager.probe_video",
                    side_effect=[low, high],
                ),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                library.archive_pptx_videos(deck, source_quality="original")
                result = library.import_external_video(
                    external, source_quality="original"
                )

            family = library.families()[0]
            self.assertEqual(result["status"], "matched")
            self.assertFalse(result["promoted"])
            self.assertEqual(len(family["variants"]), 2)
            self.assertEqual(
                library.variant_path(library.source_variant(family)).read_bytes(),
                b"low-source",
            )

            library.set_source_variant(result["variant_id"])
            self.assertEqual(
                library.variant_path(library.source_variant(family)).read_bytes(),
                b"high-source",
            )

    def test_external_video_suggestion_waits_for_manual_confirmation(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.mp4"
            candidate = root / "candidate.mp4"
            first.write_bytes(b"first")
            candidate.write_bytes(b"candidate")
            library = VideoProject.create(root / "library")
            metadata = {**no_probe(Path()), "width": 1280, "height": 720}
            source_fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            candidate_fingerprint = {
                **source_fingerprint,
                "frames": ["00000000000003ff"] * 5,
            }
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=source_fingerprint,
                ),
            ):
                library.import_external_video(first, source_quality="original")
            family = library.families()[0]
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=candidate_fingerprint,
                ),
            ):
                review = library.import_external_video(
                    candidate,
                    source_quality="original",
                    defer_suggestions=True,
                )
                self.assertEqual(review["status"], "suggested")
                self.assertEqual(len(family["variants"]), 1)
                linked = library.import_external_video(
                    candidate,
                    source_quality="original",
                    family_id=family["id"],
                    manual_confirmed=True,
                )

            self.assertEqual(linked["status"], "matched")
            self.assertEqual(len(family["variants"]), 2)
            self.assertIn(sha256_file(candidate), family["known_hashes"])

    def test_manual_pptx_override_replaces_and_remembers_low_quality_hash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            master = root / "master.pptx"
            compact = root / "compact.pptx"
            make_video_pptx(master, b"master-video", "Master")
            make_video_pptx(compact, b"compact-video", "Compact")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(master, source_quality="original")
            family = library.families()[0]
            media_path = next(iter(scan_embedded_videos(compact)))
            with ZipFile(compact) as archive:
                compact_digest = hashlib.sha256(archive.read(media_path)).hexdigest()

            result = library.upgrade_pptx_from_library(
                compact,
                family_overrides={media_path: family["id"]},
                remember_manual_matches={media_path},
            )

            self.assertEqual(result["manual_matched"], 1)
            self.assertIn(compact_digest, family["known_hashes"])
            with ZipFile(result["output_pptx"]) as archive:
                self.assertEqual(archive.read(media_path), b"master-video")

    def test_full_pptx_review_includes_exact_match_without_mutation(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            review_root = root / "review"
            make_video_pptx(source, b"source-video", "Source")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(source, source_quality="original")
                revision = library.data["revision"]
                items = library.review_pptx_matches(
                    source, review_root, include_resolved=True
                )

            self.assertEqual(len(items), 1)
            self.assertEqual(items[0]["match_kind"], "exact")
            self.assertTrue(items[0]["already_high_quality"])
            self.assertTrue(Path(items[0]["source"]).is_file())
            self.assertEqual(library.data["revision"], revision)

    def test_keep_current_media_bypasses_known_library_match(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Source")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(source, source_quality="original")
            media_path = next(iter(scan_embedded_videos(source)))

            result = library.upgrade_pptx_from_library(
                source, keep_current_media={media_path}
            )

            self.assertIsNone(result["output_pptx"])
            self.assertEqual(result["kept_current"], 1)
            self.assertEqual(result["already_high_quality"], 0)

    def test_full_review_keeps_exact_match_when_library_source_is_missing(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Source")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(source, source_quality="original")
                family = library.families()[0]
                library.require_variant_path(library.source_variant(family)).unlink()
                items = library.review_pptx_matches(
                    source, root / "review", include_resolved=True
                )

            self.assertEqual(items[0]["match_kind"], "exact")
            self.assertIsNone(items[0]["target_source"])
            self.assertTrue(items[0]["target_error"])

    def test_manual_override_can_replace_exact_match_without_learning_alias(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-video", "First")
            make_video_pptx(second, b"second-video", "Second")
            library = VideoProject.create(root / "library")
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                library.archive_pptx_videos(first, source_quality="original")
                library.archive_pptx_videos(second, source_quality="original")
            families = library.families()
            first_family = library.family_by_known_hash(
                sha256_file(first.with_suffix(".mp4"))
            )
            self.assertIsNotNone(first_family)
            target_family = next(
                family for family in families if family["id"] != first_family["id"]
            )
            media_path = next(iter(scan_embedded_videos(first)))
            first_digest = sha256_file(first.with_suffix(".mp4"))

            result = library.upgrade_pptx_from_library(
                first,
                family_overrides={media_path: target_family["id"]},
                remember_manual_matches=set(),
            )

            self.assertEqual(result["manual_matched"], 1)
            self.assertNotIn(first_digest, target_family["known_hashes"])
            with ZipFile(result["output_pptx"]) as archive:
                self.assertEqual(archive.read(media_path), b"second-video")

    def test_missing_candidate_cannot_change_source_in_memory(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck = root / "deck.pptx"
            candidate = root / "candidate.mp4"
            make_video_pptx(deck, b"source", "Demo")
            candidate.write_bytes(b"candidate")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(deck, source_quality="original")
                family = library.families()[0]
                variant = library.import_variant(
                    family["id"], candidate, verify_identity=False
                )

            source_id = family["source_variant_id"]
            known_hashes = list(family["known_hashes"])
            library.variant_path(variant).unlink()

            with self.assertRaises(FileNotFoundError):
                library.set_source_variant(variant["id"])

            self.assertEqual(family["source_variant_id"], source_id)
            self.assertEqual(family["known_hashes"], known_hashes)

    def test_external_video_creates_family_and_exact_copy_is_reused(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            external = root / "external.mp4"
            external.write_bytes(b"external-source")
            renamed = root / "renamed.mp4"
            renamed.write_bytes(external.read_bytes())
            library = VideoProject.create(root / "library")
            metadata = {**no_probe(Path()), "width": 1920, "height": 1080}
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
            ):
                created = library.import_external_video(
                    external, source_quality="original", category="Client/2026"
                )
                reused = library.import_external_video(renamed)

            self.assertEqual(created["status"], "created")
            self.assertEqual(reused["status"], "existing")
            self.assertEqual(len(library.families()), 1)
            self.assertEqual(len(library.families()[0]["variants"]), 1)
            self.assertIn(
                "media/Client/2026", library.families()[0]["variants"][0]["path"]
            )
            self.assertEqual(
                library.families()[0]["variants"][0]["origin_paths"],
                [str(external.resolve()), str(renamed.resolve())],
            )

    def test_external_video_create_rolls_back_manifest_and_file(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            external = root / "external.mp4"
            external.write_bytes(b"external-source")
            library = VideoProject.create(root / "library")
            metadata = {**no_probe(Path()), "width": 640, "height": 360}
            fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [120] * 5,
                "has_audio": False,
            }
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprint,
                ),
                patch.object(library, "save", side_effect=RuntimeError("save failed")),
                self.assertRaisesRegex(RuntimeError, "save failed"),
            ):
                library.import_external_video(external, source_quality="original")

            self.assertEqual(library.families(), [])
            self.assertEqual(list((library.root / "media").rglob("*.mp4")), [])

    def test_external_video_with_ambiguous_matches_is_not_imported(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.mp4"
            second = root / "second.mp4"
            candidate = root / "candidate.mp4"
            first.write_bytes(b"first")
            second.write_bytes(b"second")
            candidate.write_bytes(b"candidate")
            library = VideoProject.create(root / "library")
            metadata = {**no_probe(Path()), "width": 640, "height": 360}
            fingerprints = [
                {
                    "duration_ms": 2000,
                    "aspect_ppm": 1_777_778,
                    "frames": [frame] * 5,
                    "luma": [120] * 5,
                    "has_audio": False,
                }
                for frame in ("0" * 16, "f" * 16, "a" * 16)
            ]
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    side_effect=fingerprints[:2],
                ),
            ):
                library.import_external_video(first, source_quality="original")
                library.import_external_video(second, source_quality="original")

            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=fingerprints[2],
                ),
                patch.object(
                    library, "_family_fingerprint", return_value=fingerprints[2]
                ),
            ):
                result = library.import_external_video(
                    candidate, source_quality="original"
                )

            self.assertEqual(result["status"], "ambiguous")
            self.assertEqual(len(result["candidates"]), 2)
            self.assertEqual(sum(len(f["variants"]) for f in library.families()), 2)

    def test_manual_version_import_rejects_different_content(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck = root / "deck.pptx"
            external = root / "wrong.mp4"
            make_video_pptx(deck, b"source", "Demo")
            external.write_bytes(b"different")
            library = VideoProject.create(root / "library")
            source_fingerprint = {
                "duration_ms": 2000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [100] * 5,
                "has_audio": False,
            }
            different_fingerprint = {
                **source_fingerprint,
                "frames": ["f" * 16] * 5,
            }
            metadata = {**no_probe(Path()), "width": 640, "height": 360}
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=source_fingerprint,
                ),
            ):
                library.archive_pptx_videos(deck, source_quality="original")
            family = library.families()[0]
            with (
                patch("pptx_tools.video_manager.probe_video", return_value=metadata),
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=different_fingerprint,
                ),
                self.assertRaisesRegex(ValueError, "does not safely match"),
            ):
                library.import_variant(family["id"], external)

            self.assertEqual(len(family["variants"]), 1)

    def test_merge_keeps_all_known_compressed_hashes(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-source", "First")
            make_video_pptx(second, b"second-source", "Second")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(first)
                library.archive_pptx_videos(second)
            source, target = library.families()
            source_digest = library.source_variant(source)["sha256"]
            source["known_hashes"].append("compressed-alias")
            library.save()
            library.add_deck(first)
            library.add_deck(second)

            impact = library.family_merge_impact(source["id"], target["id"])
            self.assertEqual(impact["deck_count"], 1)
            self.assertEqual(impact["reference_count"], 1)
            self.assertEqual(impact["variant_count"], 1)
            self.assertIn(first.name, impact["deck_names"])

            library.merge_families(
                source["id"], target["id"], confirmed_same_content=True
            )

            self.assertEqual(len(library.families()), 1)
            self.assertIn("compressed-alias", target["known_hashes"])
            self.assertNotIn(source_digest, target["source_hashes"])
            self.assertTrue(
                all(
                    asset["family_id"] == target["id"]
                    for deck in library.decks()
                    for asset in deck["assets"]
                )
            )
            result = library.upgrade_pptx_from_library(first)
            with ZipFile(result["output_pptx"]) as archive:
                self.assertEqual(
                    archive.read(next(iter(scan_embedded_videos(first)))),
                    b"second-source",
                )

    def test_merge_rolls_back_when_manifest_save_fails(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-source", "First")
            make_video_pptx(second, b"second-source", "Second")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(first)
                library.archive_pptx_videos(second)
            source_id, target_id = [family["id"] for family in library.families()]

            with patch.object(library, "save", side_effect=RuntimeError("conflict")):
                with self.assertRaisesRegex(RuntimeError, "conflict"):
                    library.merge_families(
                        source_id,
                        target_id,
                        confirmed_same_content=True,
                    )

            self.assertEqual(
                {family["id"] for family in library.families()},
                {source_id, target_id},
            )

    def test_merge_rejects_damaged_deck_reference_without_mutation(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-source", "First")
            make_video_pptx(second, b"second-source", "Second")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(first)
                library.archive_pptx_videos(second)
            source, target = library.families()
            deck = library.add_deck(first)
            deck["assets"][0]["original_variant_id"] = "missing-variant"
            target_variant_count = len(target["variants"])

            with self.assertRaisesRegex(ValueError, "原始视频版本不存在"):
                library.merge_families(
                    source["id"], target["id"], confirmed_same_content=True
                )

            self.assertEqual(len(library.families()), 2)
            self.assertEqual(len(target["variants"]), target_variant_count)

    def test_merge_rejects_unconfirmed_different_families(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-source", "First")
            make_video_pptx(second, b"second-source", "Second")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                library.archive_pptx_videos(first)
                library.archive_pptx_videos(second)
            source_id, target_id = [family["id"] for family in library.families()]

            with self.assertRaisesRegex(ValueError, "必须人工确认"):
                library.merge_families(source_id, target_id)

            self.assertEqual(len(library.families()), 2)

    def test_content_match_requires_one_unique_conservative_candidate(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            fingerprint = {
                "duration_ms": 10_000,
                "aspect_ppm": 1_777_778,
                "frames": ["0123456789abcdef"] * 5,
                "luma": [120] * 5,
            }
            candidate = {
                "id": "first",
                "content_fingerprint": fingerprint,
                "variants": [],
            }
            library.data["families"] = [candidate]

            self.assertIs(
                library.family_by_content_fingerprint(dict(fingerprint)), candidate
            )
            macroblock_crop = {**fingerprint, "aspect_ppm": 1_791_045}
            self.assertIs(
                library.family_by_content_fingerprint(macroblock_crop), candidate
            )
            different_aspect = {**fingerprint, "aspect_ppm": 1_850_000}
            self.assertIsNone(library.family_by_content_fingerprint(different_aspect))
            different_frames = {
                **fingerprint,
                "frames": ["fedcba9876543210"] * 5,
            }
            self.assertIsNone(library.family_by_content_fingerprint(different_frames))

            audio_fingerprint = {
                **fingerprint,
                "has_audio": True,
                "audio_hash": "0011223344556677",
                "audio_luma": 80,
            }
            candidate["content_fingerprint"] = audio_fingerprint
            self.assertIs(
                library.family_by_content_fingerprint(dict(audio_fingerprint)),
                candidate,
            )
            different_audio = {
                **audio_fingerprint,
                "audio_hash": "ffeeddccbbaa9988",
            }
            self.assertIsNone(library.family_by_content_fingerprint(different_audio))
            no_audio = {
                **fingerprint,
                "has_audio": False,
            }
            self.assertIsNone(library.family_by_content_fingerprint(no_audio))
            different_length = {
                **audio_fingerprint,
                "duration_ms": 10_300,
            }
            self.assertIsNone(library.family_by_content_fingerprint(different_length))

            candidate["content_fingerprint"] = fingerprint
            library.data["families"].append(
                {
                    "id": "second",
                    "content_fingerprint": dict(fingerprint),
                    "variants": [],
                }
            )
            self.assertIsNone(library.family_by_content_fingerprint(dict(fingerprint)))

    def test_corrupted_fingerprints_degrade_to_no_match(self) -> None:
        from pptx_tools.video_manager import _fingerprints_match

        good = {
            "duration_ms": 10_000,
            "aspect_ppm": 1_777_778,
            "frames": ["0123456789abcdef"] * 5,
            "luma": [120] * 5,
            "has_audio": True,
            "audio_hash": "0011223344556677",
            "audio_luma": 80,
        }
        corrupted = [
            {},
            {"aspect_ppm": 1_777_778},
            {"duration_ms": 10_000},
            {"duration_ms": "abc", "aspect_ppm": 1_777_778},
            {"duration_ms": None, "aspect_ppm": None},
            {"duration_ms": 10_000, "aspect_ppm": 1_777_778, "frames": "oops"},
            {
                "duration_ms": 10_000,
                "aspect_ppm": 1_777_778,
                "frames": ["zz", "zz", "zz", "zz", "zz"],
                "luma": [1] * 5,
                "has_audio": False,
            },
        ]
        for broken in corrupted:
            self.assertFalse(_fingerprints_match(broken, good), broken)
            self.assertFalse(_fingerprints_match(good, broken), broken)

        library_root = Path(tempfile.mkdtemp()) / "library"
        library = VideoProject.create(library_root)
        library.data["families"] = [
            {
                "id": "broken",
                "content_fingerprint": {"unexpected": True},
                "variants": [],
            }
        ]
        self.assertIsNone(library.family_by_content_fingerprint(dict(good)))
        shutil.rmtree(library_root.parent, ignore_errors=True)

    def test_content_match_tolerates_reencoding_noise(self) -> None:
        from pptx_tools.video_manager import (
            _fingerprint_confidence,
            _fingerprints_match,
        )

        source = {
            "duration_ms": 10_000,
            "aspect_ppm": 1_777_778,
            "frames": ["0000000000000000"] * 5,
            "luma": [100] * 5,
            "has_audio": True,
            "audio_hash": "0000000000000000",
            "audio_luma": 60,
        }
        reencoded = {
            **source,
            "frames": [f"{(1 << bits) - 1:016x}" for bits in (6, 2, 3, 5, 4)],
            "audio_hash": f"{(1 << 13) - 1:016x}",
            "audio_luma": 113,
        }
        shifted_frame = {
            **source,
            "frames": [f"{(1 << bits) - 1:016x}" for bits in (1, 0, 1, 24, 0)],
            "luma": [100, 100, 100, 112, 101],
        }
        different_audio = {
            **reencoded,
            "audio_hash": f"{(1 << 14) - 1:016x}",
        }

        self.assertTrue(_fingerprints_match(source, reencoded))
        self.assertTrue(_fingerprints_match(source, shifted_frame))
        self.assertFalse(_fingerprints_match(source, different_audio))
        self.assertTrue(_fingerprint_confidence(source, reencoded)["matched"])

    def test_compressed_hash_alias_upgrades_without_tracking_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            compressed = root / "compressed.pptx"
            make_video_pptx(source, b"high-quality-source", "Source")
            make_video_pptx(compressed, b"low-quality-proxy", "Compressed")
            library = VideoProject.create(root / "library")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                archived = library.archive_pptx_videos(source)
            source_media = next(iter(archived["media_families"]))
            report = root / "report.json"
            report.write_text(
                json.dumps(
                    {
                        "videos": [
                            {
                                "media_path": source_media,
                                "output_media_path": source_media,
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            aliases = library.register_compressed_pptx_hashes(
                compressed, report, archived["media_families"]
            )
            family = library.families()[0]
            source_variant = library.source_variant(family)
            with patch.object(
                library,
                "_delivery_master",
                return_value=(
                    library.variant_path(source_variant),
                    source_variant["sha256"],
                ),
            ):
                result = library.upgrade_pptx_from_library(compressed)

            self.assertEqual(aliases, 1)
            self.assertEqual(result["matched"], 1)
            self.assertEqual(result["unmatched"], [])
            self.assertEqual(library.decks(), [])
            with ZipFile(result["output_pptx"]) as archive:
                self.assertEqual(archive.read(source_media), b"high-quality-source")

    def test_optimized_output_accepts_renamed_source_alias(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            renamed = root / "renamed.pptx"
            optimized = root / "optimized.pptx"
            make_video_pptx(source, b"source-video", "Source")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)
                shutil.copyfile(source, renamed)
                project.add_deck(renamed)
            shutil.copyfile(source, optimized)

            record = project.register_optimized_output(renamed, optimized)

            self.assertIsNotNone(record)
            self.assertEqual(len(deck["optimized_outputs"]), 1)

    def test_registered_optimized_output_reuses_deck_when_added_again(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            optimized = root / "optimized.pptx"
            renamed = root / "renamed-optimized.pptx"
            make_video_pptx(source, b"source-video", "Source")
            make_video_pptx(optimized, b"compressed-video", "Optimized")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)
                project.register_optimized_output(source, optimized)
                shutil.copy2(optimized, renamed)
                reused = project.add_deck(renamed)

            self.assertEqual(reused["id"], deck["id"])
            self.assertEqual(len(project.decks()), 1)
            self.assertIn(project.encode_path(renamed), deck["source_aliases"])

    def test_latest_optimized_output_replaces_previous_record(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(source, b"source-video", "Source")
            make_video_pptx(first, b"first-video", "First")
            make_video_pptx(second, b"second-video", "Second")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)

            first_record = project.register_optimized_output(source, first)
            second_record = project.register_optimized_output(source, second)

            self.assertIsNotNone(first_record)
            self.assertIsNotNone(second_record)
            self.assertEqual(first_record["id"], second_record["id"])
            self.assertEqual(len(deck["optimized_outputs"]), 1)
            self.assertEqual(
                project.resolve_path(second_record["path"]), second.resolve()
            )

    def test_renamed_and_moved_video_is_relinked_by_hash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck_path = root / "source.pptx"
            make_video_pptx(deck_path, b"rename-me", "Rename")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(deck_path)
            variant = project.families()[0]["variants"][0]
            renamed = project.rename_variant_file(variant["id"], "5.2追溯")
            self.assertEqual(
                renamed.name,
                f"5.2追溯_[320x180_2.0s]_{variant['sha256'][:8]}.mp4",
            )
            self.assertEqual(
                project.rename_variant_file(variant["id"], renamed.name), renamed
            )
            archive = root / "external" / "archive"
            archive.mkdir(parents=True)
            moved = archive / "final-name.mp4"
            shutil.move(renamed, moved)

            results = project.relink_missing([root / "external"])

            self.assertEqual(len(results), 1)
            self.assertEqual(project.variant_path(variant), moved.resolve())
            self.assertEqual(sha256_file(moved), variant["sha256"])

    def test_relink_refreshes_copy_changed_timestamp_after_hash_verification(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck_path = root / "source.pptx"
            make_video_pptx(deck_path, b"copied-video", "Copied")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(deck_path)
            variant = project.families()[0]["variants"][0]
            path = project.variant_path(variant)
            path.touch()

            self.assertEqual(project.status(variant), "metadata_drift")
            with self.assertRaisesRegex(ValueError, "只能隔离文件异常"):
                project.quarantine_abnormal_variant(variant["id"])
            self.assertEqual(project.relink_missing(), [])
            self.assertEqual(project.status(variant), "available")
            self.assertEqual(
                VideoProject.open(project.root).status(variant), "available"
            )

    def test_refresh_modified_variants_rebaselines_timestamp_drift(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck_path = root / "source.pptx"
            make_video_pptx(deck_path, b"drift-video", "Drift")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(deck_path)
            variant = project.families()[0]["variants"][0]
            project.variant_path(variant).touch()

            self.assertEqual(project.status(variant), "metadata_drift")
            result = project.refresh_modified_variants()

            self.assertEqual(result, {"refreshed": 1, "stale": 0})
            self.assertEqual(
                VideoProject.open(project.root).status(variant), "available"
            )

    def test_refresh_modified_variants_keeps_truly_changed_files_flagged(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            deck_path = root / "source.pptx"
            make_video_pptx(deck_path, b"changed-video", "Changed")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(deck_path)
            variant = project.families()[0]["variants"][0]
            path = project.variant_path(variant)
            with path.open("ab") as handle:
                handle.write(b"tampered")

            self.assertEqual(project.status(variant), "modified")
            result = project.refresh_modified_variants()

            self.assertEqual(result, {"refreshed": 0, "stale": 0})
            self.assertEqual(
                VideoProject.open(project.root).status(variant), "modified"
            )

    def test_detach_and_restore_keep_slide_xml_and_original_video(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"original-video-bytes", "Round trip")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)

            def placeholder(_: Path, target: Path, __: dict[str, object]) -> None:
                target.write_bytes(b"small-placeholder")

            detached = project.detach_deck(deck["id"], placeholder_builder=placeholder)
            restored = project.restore_deck(deck["id"], detached)
            original_part = deck["assets"][0]["part_path"]
            placeholder_part = deck["assets"][0]["placeholder_part"]

            with (
                ZipFile(source) as original_zip,
                ZipFile(detached) as detached_zip,
                ZipFile(restored) as restored_zip,
            ):
                slide_path = deck["assets"][0]["occurrences"][0]["slide_path"]
                self.assertEqual(
                    original_zip.read(slide_path), detached_zip.read(slide_path)
                )
                self.assertEqual(
                    original_zip.read(slide_path), restored_zip.read(slide_path)
                )
                self.assertNotIn(original_part, detached_zip.namelist())
                self.assertIn(placeholder_part, detached_zip.namelist())
                self.assertEqual(
                    restored_zip.read(original_part), b"original-video-bytes"
                )

            restored_scan = scan_embedded_videos(restored)
            self.assertIn(original_part, restored_scan)
            self.assertEqual(
                restored_scan[original_part].occurrences[0].shape_id,
                deck["assets"][0]["occurrences"][0]["shape_id"],
            )

    def test_compression_adds_and_activates_a_version(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Compress")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            original = family["variants"][0]

            def fake_compress(args, **_: object) -> dict[str, object]:
                Path(args.output).write_bytes(b"compressed-video")
                return {"output_pptx": Path(args.output)}

            with (
                patch(
                    "pptx_tools.video_manager.compact_standalone_video",
                    side_effect=fake_compress,
                ),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                compressed = project.compress_variant(original["id"], "balanced")

            self.assertEqual(len(family["variants"]), 2)
            self.assertEqual(family["active_variant_id"], compressed["id"])
            self.assertEqual(compressed["source_variant_id"], original["id"])

    def test_existing_compression_result_is_reactivated(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Compress twice")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            original = family["variants"][0]

            def fake_compress(args, **_: object) -> dict[str, object]:
                Path(args.output).write_bytes(b"stable-compressed-video")
                return {"output_pptx": Path(args.output)}

            with (
                patch(
                    "pptx_tools.video_manager.compact_standalone_video",
                    side_effect=fake_compress,
                ),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                compressed = project.compress_variant(original["id"], "balanced")
                project.activate_variant(original["id"])
                repeated = project.compress_variant(original["id"], "balanced")

            self.assertEqual(repeated["id"], compressed["id"])
            self.assertEqual(family["active_variant_id"], compressed["id"])
            self.assertEqual(len(family["variants"]), 2)

    def test_compression_can_keep_current_active_version(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Delivery")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            original = family["variants"][0]

            def fake_compress(args, **_: object) -> dict[str, object]:
                Path(args.output).write_bytes(b"delivery-video")
                return {"output_pptx": Path(args.output)}

            with (
                patch(
                    "pptx_tools.video_manager.compact_standalone_video",
                    side_effect=fake_compress,
                ),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                project.compress_variant(original["id"], "aggressive", activate=False)

            self.assertEqual(family["active_variant_id"], original["id"])
            self.assertEqual(len(family["variants"]), 2)

    def test_manifest_recovers_from_last_valid_backup(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Backup")
            project_root = root / "project"
            project = VideoProject.create(project_root)
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            project.rename_family(project.families()[0]["id"], "renamed")
            (project_root / "video-project.json").write_text(
                "{broken", encoding="utf-8"
            )

            recovered = VideoProject.open(project_root)

            self.assertTrue(recovered.recovered_from_backup)
            self.assertTrue(recovered.recovery_detail)
            self.assertEqual(len(recovered.decks()), 1)
            self.assertEqual(len(recovered.families()), 1)
            self.assertEqual(
                VideoProject.open(project_root).data["project_id"],
                recovered.data["project_id"],
            )

    def test_nested_manifest_corruption_recovers_from_backup(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Nested corruption")
            project_root = root / "project"
            project = VideoProject.create(project_root)
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            project.rename_family(project.families()[0]["id"], "backup-ready")
            manifest = project.data.copy()
            manifest["families"] = [{"id": "broken"}]
            (project_root / "video-project.json").write_text(
                json.dumps(manifest), encoding="utf-8"
            )

            recovered = VideoProject.open(project_root)

            self.assertTrue(recovered.recovered_from_backup)
            self.assertEqual(len(recovered.families()), 1)

    def test_manifest_rejects_relative_path_outside_project(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            project = VideoProject.create(Path(temp_dir) / "project")
            project.data["decks"] = [
                {
                    "id": "deck",
                    "name": "unsafe.pptx",
                    "source_path": "../unsafe.pptx",
                    "source_sha256": "digest",
                    "assets": [],
                }
            ]
            project.manifest_path.write_text(json.dumps(project.data), encoding="utf-8")

            with self.assertRaisesRegex(ValueError, "escapes"):
                VideoProject._read_manifest(project.manifest_path)

    def test_library_does_not_persist_detailed_history(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            project = VideoProject.create(Path(temp_dir) / "project")
            project.record("test_action")
            self.assertFalse((project.root / "history.jsonl").exists())

    def test_project_relative_source_survives_project_folder_move(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            project_root = root / "project"
            source = project_root / "sources" / "source.pptx"
            source.parent.mkdir(parents=True)
            make_video_pptx(source, b"portable-video", "Portable")
            project = VideoProject.create(project_root)
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)
            self.assertEqual(deck["source_path"], "sources/source.pptx")

            moved_root = root / "archive" / "moved-project"
            moved_root.parent.mkdir()
            shutil.move(project_root, moved_root)
            reopened = VideoProject.open(moved_root)

            self.assertEqual(
                reopened.deck_source_path(reopened.decks()[0]),
                (moved_root / "sources" / "source.pptx").resolve(),
            )
            self.assertEqual(reopened.source_status(reopened.decks()[0]), "available")

    def test_stale_project_instance_cannot_overwrite_newer_changes(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Concurrent")
            project_root = root / "project"
            project = VideoProject.create(project_root)
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            first = VideoProject.open(project_root)
            stale = VideoProject.open(project_root)
            family_id = first.families()[0]["id"]
            first.rename_family(family_id, "first-writer")

            with self.assertRaisesRegex(RuntimeError, "another window"):
                stale.rename_family(family_id, "stale-writer")

            self.assertEqual(
                VideoProject.open(project_root).families()[0]["name"], "first-writer"
            )

    def test_optimized_pptx_is_persisted_and_relinked_after_manual_move(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            optimized = root / "source_compacted.pptx"
            make_video_pptx(source, b"source-video", "Source")
            make_video_pptx(optimized, b"optimized-video", "Optimized")
            project_root = root / "project"
            project = VideoProject.create(project_root)
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            record = project.register_optimized_output(source, optimized)
            assert record is not None
            nested = root / "archive" / "optimized" / "renamed.pptx"
            nested.parent.mkdir(parents=True)
            shutil.move(optimized, nested)

            reopened = VideoProject.open(project_root)
            results = reopened.relink_missing_pptx([root / "archive"])

            self.assertEqual(len(results), 1)
            persisted = VideoProject.open(project_root).find_output(record["id"])[1]
            self.assertEqual(reopened.resolve_path(persisted["path"]), nested.resolve())

    def test_uppercase_pptx_is_relinked_by_hash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            optimized = root / "optimized.pptx"
            make_video_pptx(source, b"source-video", "Source")
            make_video_pptx(optimized, b"optimized-video", "Optimized")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            record = project.register_optimized_output(source, optimized)
            assert record is not None
            moved = root / "archive" / "RENAMED.PPTX"
            moved.parent.mkdir()
            shutil.move(optimized, moved)

            results = project.relink_missing_pptx([moved.parent])

            self.assertEqual(len(results), 1)
            self.assertEqual(project.resolve_path(record["path"]), moved.resolve())

    def test_variant_move_rolls_back_when_manifest_save_conflicts(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Rollback")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            variant = project.families()[0]["variants"][0]
            original = project.variant_path(variant)
            destination = root / "archive"

            with patch.object(project, "save", side_effect=RuntimeError("conflict")):
                with self.assertRaisesRegex(RuntimeError, "conflict"):
                    project.move_variant(variant["id"], destination)

            self.assertTrue(original.is_file())
            self.assertFalse((destination / original.name).exists())
            self.assertEqual(project.variant_path(variant), original)

    def test_family_move_updates_category_and_rolls_back_on_save_conflict(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Move family")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            second = root / "second.mp4"
            second.write_bytes(b"second-version")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.import_variant(family["id"], second, verify_identity=False)
            originals = {
                variant["id"]: project.variant_path(variant)
                for variant in family["variants"]
            }

            moved = project.move_family(family["id"], "示例/分类")

            self.assertEqual(
                set(moved),
                {
                    project.root / "media" / "示例" / "分类" / path.name
                    for path in originals.values()
                },
            )
            self.assertTrue(all(not path.exists() for path in originals.values()))
            self.assertTrue(all(path.is_file() for path in moved))
            self.assertEqual(family["category"], "示例/分类")
            self.assertEqual(
                {project.variant_path(variant) for variant in family["variants"]},
                set(moved),
            )

            with patch.object(
                project, "_save_locked", side_effect=RuntimeError("conflict")
            ):
                with self.assertRaisesRegex(RuntimeError, "conflict"):
                    project.move_family(family["id"], "项目/示例项目")

            self.assertTrue(all(path.is_file() for path in moved))
            self.assertEqual(family["category"], "示例/分类")
            self.assertEqual(
                {project.variant_path(variant) for variant in family["variants"]},
                set(moved),
            )

            root_paths = project.move_family(family["id"], "")
            self.assertEqual(family["category"], "")
            self.assertTrue(
                all(path.parent == project.root / "media" for path in root_paths)
            )

    def test_open_recovers_interrupted_family_move(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"source-video", "Interrupted move")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
            family = project.families()[0]
            variant = family["variants"][0]
            original = project.variant_path(variant)
            target = project.root / "media" / "示例项目" / original.name
            target.parent.mkdir(parents=True)
            shutil.move(original, target)
            (project.root / FAMILY_MOVE_JOURNAL_NAME).write_text(
                json.dumps(
                    {
                        "family_id": family["id"],
                        "moves": [
                            {
                                "variant_id": variant["id"],
                                "source": str(original),
                                "target": str(target),
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            reopened = VideoProject.open(project.root)

            self.assertTrue(original.is_file())
            self.assertFalse(target.exists())
            self.assertFalse((project.root / FAMILY_MOVE_JOURNAL_NAME).exists())
            self.assertEqual(
                reopened.variant_path(reopened.families()[0]["variants"][0]),
                original,
            )

    def test_placeholder_requires_real_timing_and_dimensions(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "unknown-video.bin"
            target = root / "placeholder.mp4"
            source.write_bytes(b"not-a-video")
            unknown = {
                "width": 0,
                "height": 0,
                "duration_sec": 0,
                "probe_error": "unavailable",
            }
            project = VideoProject.create(root / "project")

            with (
                patch("pptx_tools.video_manager.probe_video", return_value=unknown),
                patch("pptx_tools.video_manager.run_binary") as run_binary,
            ):
                with self.assertRaisesRegex(RuntimeError, "timing-safe placeholder"):
                    project.build_placeholder(source, target, unknown)

            run_binary.assert_not_called()
            self.assertFalse(target.exists())

    def test_wrong_lightweight_deck_is_rejected_without_output(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first = root / "first.pptx"
            second = root / "second.pptx"
            make_video_pptx(first, b"first-video", "First")
            make_video_pptx(second, b"second-video", "Second")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                first_deck = project.add_deck(first)
                second_deck = project.add_deck(second)

            def placeholder(_: Path, target: Path, context: dict[str, object]) -> None:
                target.write_bytes(str(context["placeholder_token"]).encode())

            project.detach_deck(first_deck["id"], placeholder_builder=placeholder)
            second_light = project.detach_deck(
                second_deck["id"], placeholder_builder=placeholder
            )
            output = root / "must-not-exist.pptx"

            with self.assertRaisesRegex(RuntimeError, "Cannot uniquely match"):
                project.restore_deck(first_deck["id"], second_light, output_path=output)

            self.assertFalse(output.exists())

    def test_custom_container_adds_content_type_when_restored(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"original-video", "Custom")
            custom = root / "custom.mov"
            custom.write_bytes(b"custom-mov-version")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                deck = project.add_deck(source)
                family = project.families()[0]
                variant = project.import_variant(
                    family["id"], custom, "mov", verify_identity=False
                )
            project.activate_variant(variant["id"])

            def placeholder(_: Path, target: Path, __: dict[str, object]) -> None:
                target.write_bytes(b"placeholder")

            detached = project.detach_deck(deck["id"], placeholder_builder=placeholder)
            restored = project.restore_deck(deck["id"], detached)
            with ZipFile(restored) as archive:
                content_types = archive.read("[Content_Types].xml")
                self.assertIn(b'Extension="mov"', content_types)
                self.assertIn(b'ContentType="video/quicktime"', content_types)

    def test_incompatible_custom_version_reports_timing_and_aspect_warnings(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"original-video", "Warnings")
            custom = root / "custom.mp4"
            custom.write_bytes(b"custom-video")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
                family = project.families()[0]
                variant = project.import_variant(
                    family["id"], custom, "custom", verify_identity=False
                )
            variant.update({"width": 320, "height": 320, "duration_sec": 8.0})
            family["content_fingerprint"] = {
                "duration_ms": 2_000,
                "aspect_ppm": 1_777_778,
                "frames": ["0" * 16] * 5,
                "luma": [100] * 5,
            }
            different_content = {
                **family["content_fingerprint"],
                "frames": ["f" * 16] * 5,
            }
            with patch(
                "pptx_tools.video_manager._video_fingerprint",
                return_value=different_content,
            ):
                warnings = project.compatibility_warnings(variant["id"])

            self.assertEqual(len(warnings), 3)
            self.assertIn("另一个视频", warnings[-1])

    def test_mp4_normalization_keeps_original_and_adds_active_derivative(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            make_video_pptx(source, b"original-video", "Normalize")
            mov = root / "legacy.mov"
            mov.write_bytes(b"legacy-video")
            project = VideoProject.create(root / "project")
            with patch("pptx_tools.video_manager.probe_video", side_effect=no_probe):
                project.add_deck(source)
                family = project.families()[0]
                legacy = project.import_variant(
                    family["id"], mov, "legacy", verify_identity=False
                )
            legacy_path = project.variant_path(legacy)

            ffmpeg_command: list[str] = []

            def fake_ffmpeg(command, **_: object) -> object:
                ffmpeg_command.extend(command)
                Path(command[-1]).write_bytes(b"normalized-mp4")
                return object()

            with (
                patch("pptx_tools.video_manager.run_binary", side_effect=fake_ffmpeg),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                compatible = project.normalize_variant_to_mp4(legacy["id"])

            self.assertTrue(legacy_path.exists())
            self.assertEqual(compatible["profile"], "mp4_high_fidelity")
            self.assertEqual(compatible["source_variant_id"], legacy["id"])
            self.assertEqual(family["active_variant_id"], compatible["id"])
            self.assertEqual(project.variant_path(compatible).suffix, ".mp4")
            self.assertEqual(ffmpeg_command[ffmpeg_command.index("-crf") + 1], "18")
            self.assertEqual(ffmpeg_command[ffmpeg_command.index("-b:a") + 1], "256k")
            self.assertNotIn("scale", " ".join(ffmpeg_command))
            self.assertNotIn("-r", ffmpeg_command)


def add_family_with_variants(
    project: VideoProject,
    name: str,
    variant_specs: list[dict[str, object]],
    *,
    known_hashes: list[str] | None = None,
    content_fingerprint: dict[str, object] | None = None,
) -> dict[str, object]:
    """Create a real family with real (fake-bytes) variant files on disk."""
    import uuid

    family_id = str(uuid.uuid4())
    variants: list[dict[str, object]] = []
    media_dir = project.root / "media" / name
    media_dir.mkdir(parents=True, exist_ok=True)
    for index, spec in enumerate(variant_specs, start=1):
        payload = f"{name}-v{index}-{spec['profile']}".encode() * max(
            1, int(spec.get("size_bytes", 100)) // 32
        )
        target = media_dir / f"{name}_v{index}.mp4"
        target.write_bytes(payload)
        digest = sha256_file(target)
        variants.append(
            {
                "id": str(uuid.uuid4()),
                "label": str(spec.get("profile", "v")),
                "profile": spec.get("profile", "custom"),
                "path": project.encode_path(target),
                "sha256": digest,
                "size_bytes": target.stat().st_size,
                "mtime_ns": target.stat().st_mtime_ns,
                "created_at": "2026-07-21T00:00:00+00:00",
                "source_variant_id": None,
                "width": spec.get("width", 640),
                "height": spec.get("height", 360),
                "duration_sec": spec.get("duration_sec", 2.0),
                "bitrate_kbps": spec.get("bitrate_kbps", 800),
                "video_codec": "h264",
                "audio_codec": "aac",
                "has_audio": spec.get("has_audio", True),
                "probe_error": "",
            }
        )
    for variant in variants[1:]:
        variant["source_variant_id"] = variants[0]["id"]
    family: dict[str, object] = {
        "id": family_id,
        "name": name,
        "category": "",
        "source_variant_id": variants[0]["id"],
        "active_variant_id": variants[0]["id"],
        "known_hashes": known_hashes or [variants[0]["sha256"]],
        "source_hashes": [variants[0]["sha256"]],
        "variants": variants,
    }
    if content_fingerprint is not None:
        family["content_fingerprint"] = content_fingerprint
    project.families().append(family)
    project.save()
    return family


FINGERPRINT_A = {
    "duration_ms": 2000,
    "aspect_ppm": 1777778,
    "frames": ["0123456789abcdef"] * 5,
    "luma": [120] * 5,
    "has_audio": True,
    "audio_hash": "0011223344556677",
    "audio_luma": 80,
}


class LibraryCleanupTest(unittest.TestCase):
    def test_focused_cleanup_only_scores_selected_family(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            selected = add_family_with_variants(
                library,
                "selected",
                [{"profile": "original"}, {"profile": "compressed"}],
            )
            add_family_with_variants(
                library,
                "other",
                [{"profile": "original"}, {"profile": "compressed"}],
            )
            with patch(
                "pptx_tools.video_manager._ssim_videos", return_value=0.99
            ) as ssim:
                groups = library.scan_cleanup_groups(focus_family_id=selected["id"])

            self.assertEqual(len(groups), 1)
            self.assertEqual(groups[0]["family_ids"], [selected["id"]])
            ssim.assert_called_once()

    def test_within_family_group_defaults_to_best_variant(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {
                        "profile": "original",
                        "width": 1920,
                        "height": 1080,
                        "bitrate_kbps": 4000,
                        "size_bytes": 10_000,
                    },
                    {
                        "profile": "high",
                        "width": 1920,
                        "height": 1080,
                        "bitrate_kbps": 1500,
                        "size_bytes": 4_000,
                    },
                    {
                        "profile": "aggressive",
                        "width": 640,
                        "height": 360,
                        "bitrate_kbps": 300,
                        "size_bytes": 1_000,
                    },
                ],
            )
            best_id = family["variants"][0]["id"]
            small_id = family["variants"][2]["id"]
            # SSIM: high->best = 0.99, aggressive->best = 0.97 (>=0.95 threshold,
            # but 360p < min(1080,1080)=1080 adequate height -> not eligible)
            with patch(
                "pptx_tools.video_manager._ssim_videos",
                side_effect=[0.99, 0.97],
            ):
                groups = library.scan_cleanup_groups(ssim_threshold=0.95)
            within = [g for g in groups if g["kind"] == "within_family"]
            self.assertEqual(len(within), 1)
            group = within[0]
            self.assertEqual(group["best_variant_id"], best_id)
            rec = group["recommendation"]
            # Cleanup defaults to the highest-quality source; the smaller
            # close version remains available as a manual alternative.
            high_id = family["variants"][1]["id"]
            self.assertEqual(rec["strategy"], "keep_best")
            self.assertEqual(rec["keep_variant_id"], best_id)
            self.assertEqual(rec["alternatives"]["keep_smallest_close"], high_id)
            # candidates carry SSIM values
            ssim_map = {c["variant_id"]: c["ssim_to_best"] for c in group["candidates"]}
            self.assertEqual(ssim_map[best_id], 1.0)
            self.assertEqual(ssim_map[small_id], 0.97)

    def test_within_family_keeps_smaller_close_as_alternative(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {
                        "profile": "original",
                        "width": 3840,
                        "height": 2160,
                        "bitrate_kbps": 20000,
                        "size_bytes": 50_000,
                    },
                    {
                        "profile": "high",
                        "width": 1920,
                        "height": 1080,
                        "bitrate_kbps": 1500,
                        "size_bytes": 4_000,
                    },
                ],
            )
            small_1080 = family["variants"][1]["id"]
            with patch("pptx_tools.video_manager._ssim_videos", side_effect=[0.97]):
                groups = library.scan_cleanup_groups(ssim_threshold=0.95)
            rec = groups[0]["recommendation"]
            # 4K->1080p SSIM 0.97 >= 0.95, height 1080 >= min(2160,1080) -> close enough
            self.assertEqual(rec["strategy"], "keep_best")
            self.assertEqual(rec["keep_variant_id"], family["variants"][0]["id"])
            self.assertEqual(rec["alternatives"]["keep_smallest_close"], small_1080)

    def test_cross_family_groups_by_shared_hash_and_fingerprint(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family_a = add_family_with_variants(
                library,
                "alpha",
                [{"profile": "original", "size_bytes": 5_000}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            # same content fingerprint -> should cluster with alpha
            family_b = add_family_with_variants(
                library,
                "beta",
                [
                    {
                        "profile": "custom",
                        "size_bytes": 2_000,
                        "width": 640,
                        "height": 360,
                    }
                ],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            # different audio fingerprint -> must NOT cluster
            other_fp = {**FINGERPRINT_A, "audio_hash": "ffeeddccbbaa9988"}
            add_family_with_variants(
                library,
                "gamma",
                [{"profile": "original", "size_bytes": 5_000}],
                content_fingerprint=other_fp,
            )
            with (
                patch(
                    "pptx_tools.video_manager._decoded_audio_correlation",
                    return_value=0.999,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.98),
            ):
                groups = library.scan_cleanup_groups()
            cross = [g for g in groups if g["kind"] == "cross_family"]
            self.assertEqual(len(cross), 1)
            self.assertEqual(
                set(cross[0]["family_ids"]), {family_a["id"], family_b["id"]}
            )

    def test_cleanup_can_group_same_named_low_bitrate_encode(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            source = {**FINGERPRINT_A, "frames": ["0" * 16] * 5}
            compressed = {
                **source,
                "frames": [f"{mask:016x}" for mask in (0x7F, 0x3F, 0x1F, 0xF, 0x7)],
            }
            first = add_family_with_variants(
                library,
                "仿真演示",
                [{"profile": "original"}],
                content_fingerprint=source,
            )
            second = add_family_with_variants(
                library,
                "仿真演示",
                [{"profile": "original"}],
                content_fingerprint=compressed,
            )
            with patch("pptx_tools.video_manager._ssim_videos", return_value=0.8):
                groups = library.scan_cleanup_groups()

            cross = [group for group in groups if group["kind"] == "cross_family"]
            self.assertEqual(len(cross), 1)
            self.assertEqual(set(cross[0]["family_ids"]), {first["id"], second["id"]})
            self.assertTrue(cross[0]["safe_to_apply"])
            self.assertTrue(
                all(item["auto_allowed"] for item in cross[0]["candidates"])
            )

    def test_cleanup_groups_copy_suffix_only_after_content_verification(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            source = {**FINGERPRINT_A, "frames": ["0" * 16] * 5}
            compressed = {
                **source,
                "frames": [f"{mask:016x}" for mask in (0x7F, 0x3F, 0x1F, 0xF, 0x7)],
            }
            first = add_family_with_variants(
                library,
                "示例视频",
                [{"profile": "original"}],
                content_fingerprint=source,
            )
            second = add_family_with_variants(
                library,
                "示例视频_1",
                [{"profile": "high"}],
                content_fingerprint=compressed,
            )

            with (
                patch(
                    "pptx_tools.video_manager._decoded_audio_correlation",
                    return_value=0.999,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.97),
            ):
                groups = library.scan_cleanup_groups()

            cross = [group for group in groups if group["kind"] == "cross_family"]
            self.assertEqual(len(cross), 1)
            self.assertEqual(set(cross[0]["family_ids"]), {first["id"], second["id"]})

    def test_cleanup_groups_same_named_ultrawide_cropped_encodes(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            source = {
                "duration_ms": 148_750,
                "aspect_ppm": 3_555_556,
                "frames": [
                    "9c1c1e0e0e0eafae",
                    "cc0f030333b3f3ff",
                    "8c2f4f39393b13cb",
                    "e421a3c3c7c14355",
                    "2f0f8b8b0e040d0d",
                ],
                "luma": [31, 23, 61, 39, 32],
                "has_audio": True,
                "audio_hash": "171b1b1313171315",
                "audio_luma": 118,
            }
            compressed = {
                **source,
                "aspect_ppm": 3_636_364,
                "frames": [
                    "9c1f1f0f0f0f8f0e",
                    "cc0f03033333f7f7",
                    "8c2f0f39393933cb",
                    "e42123c3c7c14375",
                    "09d6cb8b0e8c0e4f",
                ],
                "luma": [41, 36, 68, 50, 41],
                "audio_hash": "41001b1317171305",
            }
            first = add_family_with_variants(
                library,
                "飞书演示",
                [{"profile": "original", "width": 3840, "height": 1080}],
                content_fingerprint=source,
            )
            second = add_family_with_variants(
                library,
                "飞书演示.mp4",
                [{"profile": "high", "width": 1280, "height": 352}],
                content_fingerprint=compressed,
            )

            with (
                patch(
                    "pptx_tools.video_manager._decoded_audio_correlation",
                    return_value=0.999,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.85),
            ):
                groups = library.scan_cleanup_groups()

            cross = [group for group in groups if group["kind"] == "cross_family"]
            self.assertEqual(len(cross), 1)
            self.assertEqual(set(cross[0]["family_ids"]), {first["id"], second["id"]})
            self.assertTrue(cross[0]["safe_to_apply"])

    def test_cleanup_replaces_damaged_named_copy_with_healthy_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            healthy = add_family_with_variants(
                library,
                "demo",
                [{"profile": "original", "size_bytes": 500_000}],
            )
            damaged = add_family_with_variants(
                library,
                "demo.mp4",
                [{"profile": "original", "size_bytes": 200_000}],
            )
            healthy_variant = healthy["variants"][0]
            damaged_variant = damaged["variants"][0]
            healthy_path = library.variant_path(healthy_variant)
            damaged_path = library.variant_path(damaged_variant)
            payload = bytes(range(256)) * 2048
            healthy_path.write_bytes(payload)
            damaged_path.write_bytes(payload[:200_000])
            for variant, path in (
                (healthy_variant, healthy_path),
                (damaged_variant, damaged_path),
            ):
                variant["sha256"] = sha256_file(path)
                variant["size_bytes"] = path.stat().st_size
                variant["mtime_ns"] = path.stat().st_mtime_ns
            healthy["known_hashes"] = [healthy_variant["sha256"]]
            damaged["known_hashes"] = [damaged_variant["sha256"]]
            library.save()

            def packet_error(path: Path) -> str:
                return "partial file" if path == damaged_path else ""

            with (
                patch(
                    "pptx_tools.video_manager._video_packet_error",
                    side_effect=packet_error,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.99),
            ):
                groups = library.scan_cleanup_groups()

            cross = [group for group in groups if group["kind"] == "cross_family"]
            self.assertEqual(len(cross), 1)
            candidates = {item["variant_id"]: item for item in cross[0]["candidates"]}
            self.assertTrue(candidates[healthy_variant["id"]]["can_keep"])
            self.assertFalse(candidates[damaged_variant["id"]]["can_keep"])
            self.assertEqual(
                cross[0]["recommendation"]["keep_variant_id"],
                healthy_variant["id"],
            )

    def test_cleanup_uses_decoded_audio_for_close_named_encodes(self) -> None:
        source = {
            "duration_ms": 53_417,
            "aspect_ppm": 566_667,
            "frames": [
                "6b90b1c38714341f",
                "5b31b0c71f143c07",
                "52a1e18f1b123607",
                "a120e3850c0d0f03",
                "648347646161e0e0",
            ],
            "luma": [149, 150, 153, 146, 149],
            "has_audio": True,
            "audio_hash": "1100595919595121",
            "audio_luma": 72,
        }
        compressed = {
            **source,
            "frames": [
                "6b90b1c38714241f",
                "5b31b0c31f143c07",
                "52a1e3871b123607",
                "a120e3870d0f0f03",
                "648347646161e0e0",
            ],
            "luma": [150, 150, 152, 147, 150],
            "audio_hash": "000000000223a301",
            "audio_luma": 31,
        }
        for correlation, expected_groups in ((0.999, 1), (0.3, 0)):
            with (
                self.subTest(correlation=correlation),
                tempfile.TemporaryDirectory() as temp_dir,
            ):
                library = VideoProject.create(Path(temp_dir) / "library")
                best = add_family_with_variants(
                    library,
                    "微信视频",
                    [{"profile": "original", "bitrate_kbps": 355}],
                    content_fingerprint=source,
                )
                add_family_with_variants(
                    library,
                    "微信视频.mp4",
                    [{"profile": "high", "bitrate_kbps": 308}],
                    content_fingerprint=compressed,
                )
                with (
                    patch(
                        "pptx_tools.video_manager._decoded_audio_correlation",
                        return_value=correlation,
                    ),
                    patch("pptx_tools.video_manager._ssim_videos", return_value=0.97),
                ):
                    groups = library.scan_cleanup_groups()
                cross = [group for group in groups if group["kind"] == "cross_family"]
                self.assertEqual(len(cross), expected_groups)
                if cross:
                    self.assertEqual(
                        cross[0]["recommendation"]["keep_variant_id"],
                        best["variants"][0]["id"],
                    )

    def test_cross_family_blocks_audio_mismatch_and_trim(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            add_family_with_variants(
                library,
                "alpha",
                [{"profile": "original", "size_bytes": 5_000}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            # same frames/duration but completely different audio -> no cluster
            trimmed_fp = {
                **FINGERPRINT_A,
                "duration_ms": 1200,
                "frames": ["0123456789abcde0"] * 5,
            }
            add_family_with_variants(
                library,
                "trimmed",
                [{"profile": "custom", "size_bytes": 2_000}],
                content_fingerprint=trimmed_fp,
            )
            with patch("pptx_tools.video_manager._ssim_videos", return_value=0.98):
                groups = library.scan_cleanup_groups()
            cross = [g for g in groups if g["kind"] == "cross_family"]
            self.assertEqual(cross, [])
            self.assertEqual(len(library.families()), 2)

    def test_cross_family_fingerprint_match_requires_decoded_audio(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library,
                "alpha",
                [{"profile": "original"}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            second = add_family_with_variants(
                library,
                "beta",
                [{"profile": "original"}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            with patch(
                "pptx_tools.video_manager._decoded_audio_correlation",
                return_value=0.3,
            ):
                self.assertFalse(library._families_have_same_content(first, second))

    def test_apply_plan_quarantines_and_restores(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            keep_id = family["variants"][0]["id"]
            remove_id = family["variants"][1]["id"]
            remove_hash = family["variants"][1]["sha256"]
            remove_path = library.variant_path(family["variants"][1])
            remove_bytes = remove_path.read_bytes()
            self.assertTrue(remove_path.is_file())

            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                    }
                ]
            )
            self.assertEqual(result["failed"], 0)
            self.assertFalse(remove_path.is_file())
            pending = library.pending_cleanup()
            self.assertEqual(len(pending), 1)
            self.assertTrue(pending[0]["exists"])
            self.assertFalse(Path(pending[0]["quarantined_path"]).is_absolute())
            self.assertFalse(Path(pending[0]["original_path"]).is_absolute())
            self.assertEqual(len(library.families()[0]["variants"]), 1)
            self.assertIn(remove_hash, library.families()[0]["known_hashes"])
            # source/active still valid
            self.assertEqual(library.families()[0]["source_variant_id"], keep_id)

            reimported = Path(temp_dir) / "reimported.mp4"
            reimported.write_bytes(remove_bytes)
            self.assertEqual(
                library.import_external_video(reimported)["status"], "existing"
            )
            self.assertEqual(len(library.families()[0]["variants"]), 1)

            # restore brings it back
            restored = library.restore_cleanup_entry(pending[0]["token"])
            self.assertTrue(restored.is_file())
            self.assertEqual(len(library.families()[0]["variants"]), 2)
            self.assertEqual(library.pending_cleanup(), [])

    def test_force_cleanup_quarantines_locked_within_family_variant(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "duration_sec": 2.0},
                    {"profile": "trimmed", "duration_sec": 1.0},
                ],
            )
            keep_id = family["variants"][0]["id"]
            remove = family["variants"][1]
            remove["source_variant_id"] = None
            remove_id = remove["id"]
            remove_path = library.variant_path(remove)
            library.save()

            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                        "force_remove_variant_ids": [remove_id],
                    }
                ]
            )

            self.assertEqual(result["failed"], 0)
            self.assertFalse(remove_path.exists())
            self.assertEqual(len(library.pending_cleanup()), 1)
            self.assertEqual(
                [item["id"] for item in library.families()[0]["variants"]],
                [keep_id],
            )

    def test_force_cleanup_rejects_cross_family_decision(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library, "first", [{"profile": "original"}]
            )
            second = add_family_with_variants(library, "second", [{"profile": "copy"}])
            with self.assertRaisesRegex(ValueError, "跨族归并不支持强制整理"):
                library._validate_cleanup_decision(
                    {
                        "kind": "cross_family",
                        "keep_variant_id": first["variants"][0]["id"],
                        "remove_variant_ids": [second["variants"][0]["id"]],
                        "force_remove_variant_ids": [second["variants"][0]["id"]],
                        "merge_into_family_id": first["id"],
                        "merge_family_ids": [first["id"], second["id"]],
                    }
                )

    def test_open_recovers_interrupted_quarantine_move(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library, "demo", [{"profile": "original", "size_bytes": 10_000}]
            )
            variant = family["variants"][0]
            source = library.variant_path(variant)
            target = library.cleanup_dir / f"recovery_{source.name}"
            entry = {
                "token": "recovery-token",
                "family_id": family["id"],
                "family_name": family["name"],
                "variant": variant,
                "original_path": str(source),
                "quarantined_path": str(target),
                "quarantined_at": "2026-07-31T00:00:00+00:00",
                "reason": "test",
                "state": "moving",
            }
            library._write_cleanup_index([entry])
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(source), target)

            reopened = VideoProject.open(library.root)

            pending = reopened.pending_cleanup()
            self.assertEqual(pending[0]["state"], "quarantined")
            self.assertTrue(pending[0]["exists"])
            self.assertEqual(reopened.restore_cleanup_entry("recovery-token"), source)
            self.assertTrue(source.is_file())
            self.assertEqual(reopened.pending_cleanup(), [])

    def test_abnormal_unused_variant_can_only_be_quarantined(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            library = VideoProject.create(root / "library")
            media = library.root / "media"
            media.mkdir(exist_ok=True)
            healthy_path = media / "healthy.mp4"
            damaged_path = media / "damaged.mp4"
            healthy_path.write_bytes(b"healthy")
            damaged_path.write_bytes(b"damaged")
            healthy = {
                "id": "healthy",
                "label": "source",
                "profile": "original",
                "path": library.encode_path(healthy_path),
                "sha256": sha256_file(healthy_path),
                "size_bytes": healthy_path.stat().st_size,
                "mtime_ns": healthy_path.stat().st_mtime_ns,
                **no_probe(healthy_path),
            }
            damaged = {
                **healthy,
                "id": "damaged",
                "label": "broken",
                "path": library.encode_path(damaged_path),
                "sha256": sha256_file(damaged_path),
                "size_bytes": damaged_path.stat().st_size,
                "mtime_ns": damaged_path.stat().st_mtime_ns,
                "probe_error": "moov atom not found",
            }
            library.families().append(
                {
                    "id": "family",
                    "name": "Demo",
                    "source_variant_id": healthy["id"],
                    "active_variant_id": healthy["id"],
                    "known_hashes": [healthy["sha256"], damaged["sha256"]],
                    "source_hashes": [healthy["sha256"]],
                    "variants": [healthy, damaged],
                }
            )
            library.save()

            library.quarantine_abnormal_variant(damaged["id"])

            self.assertFalse(damaged_path.exists())
            self.assertEqual(
                [item["id"] for item in library.families()[0]["variants"]],
                ["healthy"],
            )
            self.assertEqual(len(library.pending_cleanup()), 1)
            self.assertTrue(healthy_path.is_file())

    def test_cleanup_migrates_deck_variant_references(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            first_path = root / "first.pptx"
            second_path = root / "second.pptx"
            make_video_pptx(first_path, b"first-encoding", "First")
            make_video_pptx(second_path, b"second-encoding", "Second")
            library = VideoProject.create(root / "library")
            with (
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
                patch("pptx_tools.video_manager._video_fingerprint", return_value=None),
            ):
                first_deck = library.add_deck(first_path)
                second_deck = library.add_deck(second_path)
            first, second = library.families()
            first["content_fingerprint"] = dict(FINGERPRINT_A)
            second["content_fingerprint"] = dict(FINGERPRINT_A)
            library.save()
            keep_id = first["variants"][0]["id"]
            remove_id = second["variants"][0]["id"]

            with patch(
                "pptx_tools.video_manager._decoded_audio_correlation",
                return_value=0.999,
            ):
                result = library.apply_cleanup_plan(
                    [
                        {
                            "kind": "cross_family",
                            "keep_variant_id": keep_id,
                            "remove_variant_ids": [remove_id],
                            "merge_into_family_id": first["id"],
                            "merge_family_ids": [first["id"], second["id"]],
                        }
                    ]
                )

            self.assertEqual(result["failed"], 0)
            for deck in (first_deck, second_deck):
                self.assertEqual(deck["assets"][0]["family_id"], first["id"])
                self.assertEqual(deck["assets"][0]["original_variant_id"], keep_id)

    def test_apply_plan_rolls_back_on_save_failure(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            keep_id = family["variants"][0]["id"]
            remove_id = family["variants"][1]["id"]
            remove_path = library.variant_path(family["variants"][1])
            with patch.object(
                VideoProject, "save", side_effect=RuntimeError("conflict")
            ):
                result = library.apply_cleanup_plan(
                    [
                        {
                            "kind": "within_family",
                            "family_id": family["id"],
                            "keep_variant_id": keep_id,
                            "remove_variant_ids": [remove_id],
                        }
                    ]
                )
            self.assertEqual(result["failed"], 1)
            self.assertTrue(remove_path.is_file())
            self.assertEqual(len(library.families()[0]["variants"]), 2)
            self.assertEqual(library.pending_cleanup(), [])

    def test_empty_cleanup_blocked_until_migration_complete(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            keep_id = family["variants"][0]["id"]
            remove_id = family["variants"][1]["id"]
            library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                    }
                ]
            )
            self.assertEqual(library.cleanup_pending_issues(), [])
            removed = library.empty_cleanup()
            self.assertEqual(removed, 1)
            self.assertEqual(library.pending_cleanup(), [])

            # If a manifest record still references the quarantined variant,
            # emptying must refuse.
            family2 = add_family_with_variants(
                library,
                "demo2",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family2["id"],
                        "keep_variant_id": family2["variants"][0]["id"],
                        "remove_variant_ids": [family2["variants"][1]["id"]],
                    }
                ]
            )
            # simulate an inconsistent manifest: re-add the quarantined variant
            quarantined_variant = library._read_cleanup_index()[0]["variant"]
            library.families()[0]["variants"].append(quarantined_variant)
            library.save()
            issues = library.cleanup_pending_issues()
            self.assertTrue(issues)
            with self.assertRaises(RuntimeError):
                library.empty_cleanup()

    def test_within_family_blocks_unverified_custom_video(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "custom", "size_bytes": 4_000},
                ],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            family["variants"][1]["source_variant_id"] = None
            different = {**FINGERPRINT_A, "duration_ms": 5000}
            with (
                patch(
                    "pptx_tools.video_manager._video_fingerprint",
                    return_value=different,
                ),
                patch("pptx_tools.video_manager._ssim_videos") as ssim,
            ):
                groups = library.scan_cleanup_groups()
            within = next(group for group in groups if group["kind"] == "within_family")
            self.assertFalse(within["safe_to_apply"])
            self.assertFalse(within["candidates"][1]["auto_allowed"])
            ssim.assert_not_called()

    def test_cross_family_group_suppresses_overlapping_within_groups(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library,
                "alpha",
                [
                    {"profile": "original", "size_bytes": 5_000},
                    {"profile": "high", "size_bytes": 3_000},
                ],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            second = add_family_with_variants(
                library,
                "beta",
                [
                    {"profile": "original", "size_bytes": 5_000},
                    {"profile": "high", "size_bytes": 3_000},
                ],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            with (
                patch(
                    "pptx_tools.video_manager._decoded_audio_correlation",
                    return_value=0.999,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.98),
            ):
                groups = library.scan_cleanup_groups()
            self.assertEqual([group["kind"] for group in groups], ["cross_family"])
            self.assertEqual(set(groups[0]["family_ids"]), {first["id"], second["id"]})

    def test_apply_plan_rejects_forged_cross_family_merge(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library, "alpha", [{"profile": "original"}]
            )
            second = add_family_with_variants(
                library, "beta", [{"profile": "original"}]
            )
            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "cross_family",
                        "keep_variant_id": first["variants"][0]["id"],
                        "remove_variant_ids": [second["variants"][0]["id"]],
                        "merge_into_family_id": first["id"],
                        "merge_family_ids": [first["id"], second["id"]],
                    }
                ]
            )
            self.assertEqual(result["failed"], 1)
            self.assertEqual(len(library.families()), 2)

    def test_cross_family_save_failure_rolls_back_manifest_and_files(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library,
                "alpha",
                [{"profile": "original"}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            second = add_family_with_variants(
                library,
                "beta",
                [{"profile": "original"}],
                content_fingerprint=dict(FINGERPRINT_A),
            )
            second_path = library.variant_path(second["variants"][0])
            manifest_before = library.manifest_path.read_bytes()
            with patch.object(library, "save", side_effect=RuntimeError("conflict")):
                result = library.apply_cleanup_plan(
                    [
                        {
                            "kind": "cross_family",
                            "keep_variant_id": first["variants"][0]["id"],
                            "remove_variant_ids": [second["variants"][0]["id"]],
                            "merge_into_family_id": first["id"],
                            "merge_family_ids": [first["id"], second["id"]],
                        }
                    ]
                )
            self.assertEqual(result["failed"], 1)
            self.assertEqual(library.manifest_path.read_bytes(), manifest_before)
            self.assertTrue(second_path.is_file())
            self.assertEqual(library.pending_cleanup(), [])
            self.assertEqual(len(library.families()), 2)

    def test_cross_family_exact_duplicate_file_is_quarantined(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            first = add_family_with_variants(
                library, "alpha", [{"profile": "original"}]
            )
            second = add_family_with_variants(
                library, "beta", [{"profile": "original"}]
            )
            first_variant = first["variants"][0]
            second_variant = second["variants"][0]
            second_path = library.variant_path(second_variant)
            shutil.copyfile(library.variant_path(first_variant), second_path)
            second_variant["sha256"] = first_variant["sha256"]
            second_variant["size_bytes"] = second_path.stat().st_size
            second["known_hashes"] = [first_variant["sha256"]]
            second["source_hashes"] = [first_variant["sha256"]]
            library.save()

            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "cross_family",
                        "keep_variant_id": first_variant["id"],
                        "remove_variant_ids": [second_variant["id"]],
                        "merge_into_family_id": first["id"],
                        "merge_family_ids": [first["id"], second["id"]],
                    }
                ]
            )
            self.assertEqual(result["failed"], 0)
            self.assertFalse(second_path.exists())
            self.assertEqual(len(library.families()), 1)
            pending = library.pending_cleanup()
            self.assertEqual(len(pending), 1)
            restored = library.restore_cleanup_entry(pending[0]["token"])
            self.assertTrue(restored.is_file())
            self.assertEqual(len(library.families()[0]["variants"]), 2)
            self.assertEqual(library.pending_cleanup(), [])

    def test_cleanup_index_corruption_and_escape_fail_closed(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            library.cleanup_dir.mkdir()
            unknown = library.cleanup_dir / "unknown.mp4"
            unknown.write_bytes(b"keep")
            library.cleanup_index_path.write_text("{broken", encoding="utf-8")
            with self.assertRaises(RuntimeError):
                library.empty_cleanup()
            self.assertTrue(unknown.is_file())

            outside = Path(temp_dir) / "outside.mp4"
            outside.write_bytes(b"keep-outside")
            library.cleanup_index_path.write_text(
                json.dumps([{"quarantined_path": str(outside)}]), encoding="utf-8"
            )
            with self.assertRaises(RuntimeError):
                library.empty_cleanup()
            self.assertTrue(outside.is_file())

            held = library.cleanup_dir / "held.mp4"
            held.write_bytes(b"held")
            library.cleanup_index_path.write_text(
                json.dumps(
                    [
                        {
                            "token": "escape",
                            "original_path": str(outside),
                            "quarantined_path": library.encode_path(held),
                        }
                    ]
                ),
                encoding="utf-8",
            )
            with self.assertRaises(RuntimeError):
                library.restore_cleanup_entry("escape")
            self.assertEqual(outside.read_bytes(), b"keep-outside")

    def test_create_unified_version_uses_best_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {
                        "profile": "original",
                        "width": 3840,
                        "height": 2160,
                        "size_bytes": 50_000,
                    },
                    {
                        "profile": "aggressive",
                        "width": 640,
                        "height": 360,
                        "size_bytes": 1_000,
                    },
                ],
            )
            best_id = family["variants"][0]["id"]
            ffmpeg_calls: list[list[str]] = []

            def fake_transcode(source, target, width, height, *, family_id=""):
                ffmpeg_calls.append([str(source), str(target), family_id])
                Path(target).write_bytes(b"unified-1080p-content")

            with (
                patch(
                    "pptx_tools.video_manager._transcode_high_quality_mp4",
                    side_effect=fake_transcode,
                ),
                patch("pptx_tools.video_manager.probe_video", side_effect=no_probe),
            ):
                unified = library.create_unified_version(family["id"])
            self.assertEqual(ffmpeg_calls[0][2], family["id"])
            # encoded from the best (4K original), not the compressed variant
            self.assertIn("demo_v1.mp4", ffmpeg_calls[0][0])
            self.assertEqual(unified["source_variant_id"], best_id)
            self.assertEqual(library.families()[0]["source_variant_id"], unified["id"])

    def test_integrity_cache_survives_transitive_cross_family_merge(self) -> None:
        """Three same-named families: A~B share a hash, B~C are damaged copies.
        The damaged-copy check for C must still populate integrity_errors even
        though find(B)==find(C) skips the named-grouping step, so C's candidate
        carries a non-empty integrity_error and can_keep=False."""
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            healthy_a = add_family_with_variants(
                library,
                "demo",
                [{"profile": "original", "size_bytes": 500_000}],
            )
            # B shares A's hash -> transitive union via exact hash
            healthy_b = add_family_with_variants(
                library,
                "demo.mp4",
                [{"profile": "original", "size_bytes": 500_000}],
            )
            # Make B's variant byte-identical to A's so they share a hash
            a_path = library.variant_path(healthy_a["variants"][0])
            b_path = library.variant_path(healthy_b["variants"][0])
            payload = bytes(range(256)) * 2048
            b_path.write_bytes(payload)
            a_path.write_bytes(payload)
            for fam, variant in (
                (healthy_a, healthy_a["variants"][0]),
                (healthy_b, healthy_b["variants"][0]),
            ):
                variant["sha256"] = sha256_file(library.variant_path(variant))
                variant["size_bytes"] = library.variant_path(variant).stat().st_size
                variant["mtime_ns"] = library.variant_path(variant).stat().st_mtime_ns
            healthy_a["known_hashes"] = [healthy_a["variants"][0]["sha256"]]
            healthy_b["known_hashes"] = [healthy_b["variants"][0]["sha256"]]

            # C is a damaged (truncated) copy of the same name
            damaged_c = add_family_with_variants(
                library,
                "demo.mkv",
                [{"profile": "original", "size_bytes": 200_000}],
            )
            c_path = library.variant_path(damaged_c["variants"][0])
            c_path.write_bytes(payload[:200_000])
            c_variant = damaged_c["variants"][0]
            c_variant["sha256"] = sha256_file(c_path)
            c_variant["size_bytes"] = c_path.stat().st_size
            c_variant["mtime_ns"] = c_path.stat().st_mtime_ns
            damaged_c["known_hashes"] = [c_variant["sha256"]]
            library.save()

            def packet_error(path: Path) -> str:
                return "partial file" if path == c_path else ""

            with (
                patch(
                    "pptx_tools.video_manager._video_packet_error",
                    side_effect=packet_error,
                ),
                patch("pptx_tools.video_manager._ssim_videos", return_value=0.99),
            ):
                groups = library.scan_cleanup_groups()

            cross = [g for g in groups if g["kind"] == "cross_family"]
            # C must appear in a cross-family group alongside A/B
            all_variant_ids = {c["variant_id"] for g in cross for c in g["candidates"]}
            self.assertIn(c_variant["id"], all_variant_ids)
            # Find C's candidate and verify integrity_error propagated
            c_candidate = None
            for group in cross:
                for candidate in group["candidates"]:
                    if candidate["variant_id"] == c_variant["id"]:
                        c_candidate = candidate
                        break
                if c_candidate:
                    break
            self.assertIsNotNone(c_candidate)
            self.assertTrue(c_candidate["integrity_error"])
            self.assertFalse(c_candidate["can_keep"])

    def test_cleanup_issues_self_heal_after_restore_crash_window(self) -> None:
        """Simulate the crash window: restore committed the variant to the
        manifest (save done) but the cleanup index still has the entry.
        cleanup_pending_issues must NOT report a false 'file missing' alarm."""
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            keep_id = family["variants"][0]["id"]
            remove_variant = family["variants"][1]
            remove_id = remove_variant["id"]

            # Apply cleanup to quarantine the high variant
            library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [remove_id],
                    }
                ]
            )
            pending = library.pending_cleanup()
            self.assertEqual(len(pending), 1)
            token = pending[0]["token"]
            quarantined_path = library._cleanup_entry_path(pending[0])

            # Restore normally (this removes the index entry)
            library.restore_cleanup_entry(token)
            self.assertEqual(library.pending_cleanup(), [])

            # Now simulate the crash window: re-add the index entry as if
            # save() succeeded but _write_cleanup_index was interrupted.
            # The variant is already back in the family; the quarantined file
            # is gone (moved back to media/).
            entry = {
                "token": token,
                "family_id": family["id"],
                "family_name": family["name"],
                "variant": copy.deepcopy(remove_variant),
                "original_path": library.encode_path(
                    library.variant_path(remove_variant)
                ),
                "quarantined_path": library.encode_path(quarantined_path),
                "quarantined_at": "2026-08-05T00:00:00+00:00",
                "reason": "整理视频库",
                "state": "quarantined",
            }
            library._write_cleanup_index([entry])

            # cleanup_pending_issues must self-heal: the variant is back in the
            # family with a matching file, so no false "file missing" alarm.
            issues = library.cleanup_pending_issues()
            self.assertEqual(issues, [])

    def test_apply_cleanup_plan_skips_noop_decision(self) -> None:
        """A within_family decision with remove_variant_ids=[] (or only the
        keep_id) is a no-op: must not save()/record(), and returns skipped."""
        with tempfile.TemporaryDirectory() as temp_dir:
            library = VideoProject.create(Path(temp_dir) / "library")
            family = add_family_with_variants(
                library,
                "demo",
                [
                    {"profile": "original", "size_bytes": 10_000},
                    {"profile": "high", "size_bytes": 4_000},
                ],
            )
            keep_id = family["variants"][0]["id"]
            revision_before = library.data.get("revision", 0)

            result = library.apply_cleanup_plan(
                [
                    {
                        "kind": "within_family",
                        "family_id": family["id"],
                        "keep_variant_id": keep_id,
                        "remove_variant_ids": [],
                    }
                ]
            )

            self.assertEqual(result["applied"], 1)
            self.assertEqual(result["failed"], 0)
            self.assertTrue(result["results"][0]["skipped"])
            # No manifest revision consumed
            self.assertEqual(library.data.get("revision", 0), revision_before)
            # No audit record for a skipped decision
            self.assertFalse(
                any(
                    e.get("event") == "cleanup_plan_applied"
                    for e in library.data.get("history", [])
                )
            )
            # Variants untouched
            self.assertEqual(len(library.families()[0]["variants"]), 2)


if __name__ == "__main__":
    unittest.main()
